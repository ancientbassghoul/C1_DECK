from pathlib import Path
import argparse
import os
import subprocess
from zipfile import ZIP_DEFLATED, ZipFile

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


SCRIPT_DIR = Path(__file__).resolve().parent
DEFAULT_VISUALS_DIR = SCRIPT_DIR / "visuals"
DEFAULT_DATASET_DIR = SCRIPT_DIR / "dataset"
DEFAULT_MASCOT_PATH = (
    SCRIPT_DIR / "assets" / "Time_Traveler_Spoilers_Flattened_Transparent.png"
)
DEFAULT_SERIOUS_MASCOT_PATH = (
    SCRIPT_DIR / "assets" / "Time_Traveler_Spoilers_Red_Portal_Transparent.png"
)
ANIMATION_SCRIPT = SCRIPT_DIR / "apply_speech_animations.ps1"
OUTPUT_PPTX = SCRIPT_DIR / "presentation.pptx"

SPEECH_OVERLAY_SLIDES = (1, 4, 5, 10, 11, 14, 25)
SPEECH_TEXT_MARKERS = {
    1: "≤10 pixels",
    4: "First commit after",
    5: "You have no idea how important",
    10: "This exact commit also complains",
    11: "Yeah, so much effort",
    14: "I know the guy",
    25: "...and to presuade Claude",
}
MASCOT_X = 11.75
MASCOT_Y = 5.85
MASCOT_SIZE = 1.48
SPEECH_BUBBLE_RIGHT = 11.86
SPEECH_TAIL_TARGET_Y = MASCOT_Y + 0.08


def find_asset(assets_dir, *names):
    """Return the first matching asset from a project-local asset directory."""
    for name in names:
        exact = assets_dir / name
        if exact.exists():
            return exact
        matches = sorted(assets_dir.glob(f"*{name}"))
        if matches:
            return matches[0]
    raise FileNotFoundError(
        f"Could not find any of {names!r} in {assets_dir.resolve()}."
    )

SLIDE_W = 13.333
SLIDE_H = 7.5

PRIMARY = "101827"
SECONDARY = "27E38D"
ACCENT = "FF5B5B"
OFFWHITE = "F5F2EA"
WHITE = "FFFFFF"
MUTED = "8A94A4"
PALE = "DCE2EA"
DEEP = "09111E"
SOFT_GREEN = "DDF8EC"

HEADER_FONT = "Bahnschrift SemiBold"
BODY_FONT = "Aptos"

ACT_LABELS = [
    "Setup",
    "Bring-Up",
    "By Hand",
    "Check-In",
    "Pivot",
    "New Brain",
    "Bug Hunt",
    "Tuning",
    "Hunger",
    "Landing",
]

ACT_SPECS = [
    {"number": "0", "label": "Setup", "header": "THE SETUP", "start": 1, "end": 3},
    {"number": "I", "label": "Bring-Up", "header": "BRING-UP", "start": 4, "end": 6},
    {
        "number": "II",
        "label": "By Hand",
        "header": "BY HAND,\nFOR NOW",
        "start": 7,
        "end": 11,
    },
    {
        "number": "III",
        "label": "Check-In",
        "header": "ASKING FOR\nDIRECTIONS",
        "start": 12,
        "end": 18,
    },
    {
        "number": "IV",
        "label": "Pivot",
        "header": "STARTING OVER",
        "start": 19,
        "end": 23,
    },
    {
        "number": "V",
        "label": "New Brain",
        "header": "HOW THE NEW\nBRAIN THINKS",
        "start": 24,
        "end": 25,
    },
    {
        "number": "VI",
        "label": "Bug Hunt",
        "header": "THE BUG HUNT",
        "start": 26,
        "end": 28,
    },
    {
        "number": "VII",
        "label": "Tuning",
        "header": "TUNING THE\nMACHINE",
        "start": 29,
        "end": 32,
    },
    {
        "number": "VIII",
        "label": "Hunger",
        "header": "THE HUNGER",
        "start": 33,
        "end": 33,
    },
    {
        "number": "IX",
        "label": "Landing",
        "header": "WHERE IT LANDED",
        "start": 34,
        "end": 35,
    },
]

ACT_IV_STATIONS = (19, 20, 21, 23)
ACT_V_STATIONS = (23.5, 24, 25)


def rgb(hex_color):
    return RGBColor.from_string(hex_color)


def set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_rect(slide, x, y, w, h, fill, line=None, radius=True, line_width=1):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    font_size,
    color,
    font=BODY_FONT,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.0,
    italic=False,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return box


def add_rich_text(
    slide,
    segments,
    x,
    y,
    w,
    h,
    font_size,
    color,
    valign=MSO_ANCHOR.TOP,
    margin=0.0,
):
    """Add one text item with per-segment font choices."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    paragraph = frame.paragraphs[0]
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    for text, font in segments:
        run = paragraph.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(font_size)
        run.font.color.rgb = rgb(color)
    return box


def add_title(slide, title, dark=False, font_size=40):
    color = WHITE if dark else PRIMARY
    add_text(
        slide,
        title,
        1.75,
        0.36,
        11.1,
        0.62,
        font_size,
        color,
        font=HEADER_FONT,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.75), Inches(1.10), Inches(0.78), Inches(0.055)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(SECONDARY)
    line.line.fill.background()


def add_act_rail(slide, current_slide):
    rail = add_rect(slide, 0, 0, 1.42, SLIDE_H, PRIMARY, None, radius=False)

    add_text(
        slide,
        "ACT 0 —\nTHE SETUP",
        0.18,
        0.18,
        1.08,
        0.78,
        13.5,
        WHITE,
        font=HEADER_FONT,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.18), Inches(1.02), Inches(0.48), Inches(0.045)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(SECONDARY)
    accent.line.fill.background()

    current_y = 1.17
    add_rect(slide, 0.10, current_y, 1.20, 1.08, DEEP, SECONDARY, radius=True)
    add_text(
        slide,
        "Setup",
        0.55,
        current_y + 0.37,
        0.66,
        0.28,
        10.5,
        SECONDARY,
        font=HEADER_FONT,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )

    dot_x = 0.32
    dot_ys = [current_y + 0.24, current_y + 0.53, current_y + 0.82]
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(dot_x - 0.018),
        Inches(dot_ys[0]),
        Inches(0.036),
        Inches(dot_ys[-1] - dot_ys[0]),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(MUTED)
    line.line.fill.background()

    for index, dot_y in enumerate(dot_ys, start=1):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(dot_x - 0.07),
            Inches(dot_y - 0.07),
            Inches(0.14),
            Inches(0.14),
        )
        if index < current_slide:
            dot.fill.solid()
            dot.fill.fore_color.rgb = rgb(MUTED)
            dot.line.color.rgb = rgb(MUTED)
        elif index == current_slide:
            dot.fill.solid()
            dot.fill.fore_color.rgb = rgb(SECONDARY)
            dot.line.color.rgb = rgb(SECONDARY)
        else:
            dot.fill.solid()
            dot.fill.fore_color.rgb = rgb(PRIMARY)
            dot.line.color.rgb = rgb(PALE)
        dot.line.width = Pt(1.25)

    row_y = 2.39
    row_h = 0.43
    for label in ACT_LABELS[1:]:
        tick = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.25),
            Inches(row_y + 0.19),
            Inches(0.12),
            Inches(0.018),
        )
        tick.fill.solid()
        tick.fill.fore_color.rgb = rgb(MUTED)
        tick.line.fill.background()
        add_text(
            slide,
            label,
            0.47,
            row_y + 0.06,
            0.78,
            0.28,
            8.3,
            MUTED,
            font=BODY_FONT,
            valign=MSO_ANCHOR.MIDDLE,
        )
        row_y += row_h

    return rail


def add_act_i_rail(slide, current_slide):
    """Add the persistent rail for Act I without altering the Act 0 rail."""
    rail = add_rect(slide, 0, 0, 1.42, SLIDE_H, PRIMARY, None, radius=False)

    add_text(
        slide,
        "ACT I —\nBRING-UP",
        0.18,
        0.18,
        1.08,
        0.78,
        13.5,
        WHITE,
        font=HEADER_FONT,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.18), Inches(1.02), Inches(0.48), Inches(0.045)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(SECONDARY)
    accent.line.fill.background()

    current_y = 1.17
    add_rect(slide, 0.10, current_y, 1.20, 1.08, DEEP, SECONDARY, radius=True)
    add_text(
        slide,
        "Bring-Up",
        0.55,
        current_y + 0.37,
        0.66,
        0.28,
        10.5,
        SECONDARY,
        font=HEADER_FONT,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )

    dot_x = 0.32
    dot_ys = [current_y + 0.24, current_y + 0.53, current_y + 0.82]
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(dot_x - 0.018),
        Inches(dot_ys[0]),
        Inches(0.036),
        Inches(dot_ys[-1] - dot_ys[0]),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(MUTED)
    line.line.fill.background()

    for index, dot_y in enumerate(dot_ys, start=4):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(dot_x - 0.07),
            Inches(dot_y - 0.07),
            Inches(0.14),
            Inches(0.14),
        )
        if index < current_slide:
            dot.fill.solid()
            dot.fill.fore_color.rgb = rgb(MUTED)
            dot.line.color.rgb = rgb(MUTED)
        elif index == current_slide:
            dot.fill.solid()
            dot.fill.fore_color.rgb = rgb(SECONDARY)
            dot.line.color.rgb = rgb(SECONDARY)
        else:
            dot.fill.solid()
            dot.fill.fore_color.rgb = rgb(PRIMARY)
            dot.line.color.rgb = rgb(PALE)
        dot.line.width = Pt(1.25)

    other_acts = [(0, ACT_LABELS[0])] + list(enumerate(ACT_LABELS[2:], start=2))
    row_y = 2.39
    row_h = 0.43
    for act_index, label in other_acts:
        tick = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.25),
            Inches(row_y + 0.17),
            Inches(0.12),
            Inches(0.08),
        )
        if act_index < 1:
            tick.fill.solid()
            tick.fill.fore_color.rgb = rgb(MUTED)
            tick.line.color.rgb = rgb(MUTED)
        else:
            tick.fill.solid()
            tick.fill.fore_color.rgb = rgb(PRIMARY)
            tick.line.color.rgb = rgb(MUTED)
        tick.line.width = Pt(0.8)
        add_text(
            slide,
            label,
            0.47,
            row_y + 0.06,
            0.78,
            0.28,
            8.3,
            MUTED,
            font=BODY_FONT,
            valign=MSO_ANCHOR.MIDDLE,
        )
        row_y += row_h

    return rail


def add_act_ii_rail(slide, current_slide):
    """Add the persistent rail and five stations for Act II."""
    rail = add_rect(slide, 0, 0, 1.42, SLIDE_H, PRIMARY, None, radius=False)

    add_text(
        slide,
        "ACT II —\nBY HAND,\nFOR NOW",
        0.18,
        0.13,
        1.08,
        0.88,
        11.8,
        WHITE,
        font=HEADER_FONT,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.18), Inches(1.02), Inches(0.48), Inches(0.045)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(SECONDARY)
    accent.line.fill.background()

    current_y = 1.17
    add_rect(slide, 0.10, current_y, 1.20, 1.48, DEEP, SECONDARY, radius=True)
    add_text(
        slide,
        "By Hand",
        0.55,
        current_y + 0.57,
        0.66,
        0.30,
        10.5,
        SECONDARY,
        font=HEADER_FONT,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )

    dot_x = 0.32
    dot_ys = [current_y + 0.24 + 0.25 * index for index in range(5)]
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(dot_x - 0.018),
        Inches(dot_ys[0]),
        Inches(0.036),
        Inches(dot_ys[-1] - dot_ys[0]),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(MUTED)
    line.line.fill.background()

    for index, dot_y in enumerate(dot_ys, start=7):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(dot_x - 0.07),
            Inches(dot_y - 0.07),
            Inches(0.14),
            Inches(0.14),
        )
        if index < current_slide:
            dot.fill.solid()
            dot.fill.fore_color.rgb = rgb(MUTED)
            dot.line.color.rgb = rgb(MUTED)
        elif index == current_slide:
            dot.fill.solid()
            dot.fill.fore_color.rgb = rgb(SECONDARY)
            dot.line.color.rgb = rgb(SECONDARY)
        else:
            dot.fill.solid()
            dot.fill.fore_color.rgb = rgb(PRIMARY)
            dot.line.color.rgb = rgb(PALE)
        dot.line.width = Pt(1.25)

    other_acts = [
        (act_index, label)
        for act_index, label in enumerate(ACT_LABELS)
        if act_index != 2
    ]
    row_y = 2.78
    row_h = 0.40
    for act_index, label in other_acts:
        tick = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.25),
            Inches(row_y + 0.16),
            Inches(0.12),
            Inches(0.08),
        )
        if act_index < 2:
            tick.fill.solid()
            tick.fill.fore_color.rgb = rgb(MUTED)
            tick.line.color.rgb = rgb(MUTED)
        else:
            tick.fill.solid()
            tick.fill.fore_color.rgb = rgb(PRIMARY)
            tick.line.color.rgb = rgb(MUTED)
        tick.line.width = Pt(0.8)
        add_text(
            slide,
            label,
            0.47,
            row_y + 0.05,
            0.78,
            0.28,
            8.1,
            MUTED,
            font=BODY_FONT,
            valign=MSO_ANCHOR.MIDDLE,
        )
        row_y += row_h

    return rail


def get_act_spec(slide_number):
    if slide_number == 23.5:
        return 5, ACT_SPECS[5]
    for act_index, spec in enumerate(ACT_SPECS):
        if spec["start"] <= slide_number <= spec["end"]:
            return act_index, spec
    raise ValueError(f"Slide {slide_number} is outside the configured 35-slide deck.")


def add_inline_act_rail(slide, current_slide, station_slides=None):
    """Render the active act as an inline accordion in chronological order."""
    active_index, active_spec = get_act_spec(current_slide)
    if station_slides is None:
        station_slides = tuple(
            range(active_spec["start"], active_spec["end"] + 1)
        )
    else:
        station_slides = tuple(station_slides)
    if current_slide not in station_slides:
        raise ValueError(
            f"Current slide {current_slide} is not in the active station list."
        )
    rail = add_rect(slide, 0, 0, 1.42, SLIDE_H, PRIMARY, None, radius=False)

    header_text = f'ACT {active_spec["number"]} —\n{active_spec["header"]}'
    header_lines = header_text.count("\n") + 1
    header_font_size = 13.0 if header_lines == 2 else 11.4
    add_text(
        slide,
        header_text,
        0.18,
        0.12,
        1.08,
        0.90,
        header_font_size,
        WHITE,
        font=HEADER_FONT,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.18), Inches(1.02), Inches(0.48), Inches(0.045)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(SECONDARY)
    accent.line.fill.background()

    list_y = 1.18
    collapsed_h = 0.37
    station_step = 0.20
    station_count = len(station_slides)
    active_h = collapsed_h + station_count * station_step + 0.13

    for act_index, spec in enumerate(ACT_SPECS):
        if act_index != active_index:
            marker = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.25),
                Inches(list_y + 0.145),
                Inches(0.12),
                Inches(0.08),
            )
            if act_index < active_index:
                marker.fill.solid()
                marker.fill.fore_color.rgb = rgb(MUTED)
                marker.line.color.rgb = rgb(MUTED)
            else:
                marker.fill.solid()
                marker.fill.fore_color.rgb = rgb(PRIMARY)
                marker.line.color.rgb = rgb(MUTED)
            marker.line.width = Pt(0.8)
            add_text(
                slide,
                spec["label"],
                0.47,
                list_y + 0.045,
                0.78,
                0.28,
                8.2,
                MUTED,
                font=BODY_FONT,
                valign=MSO_ANCHOR.MIDDLE,
            )
            list_y += collapsed_h
            continue

        add_rect(
            slide,
            0.10,
            list_y,
            1.20,
            active_h,
            DEEP,
            SECONDARY,
            radius=True,
            line_width=1.0,
        )
        active_marker = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.25),
            Inches(list_y + 0.145),
            Inches(0.12),
            Inches(0.08),
        )
        active_marker.fill.solid()
        active_marker.fill.fore_color.rgb = rgb(SECONDARY)
        active_marker.line.color.rgb = rgb(SECONDARY)
        add_text(
            slide,
            spec["label"],
            0.47,
            list_y + 0.045,
            0.78,
            0.28,
            9.0,
            SECONDARY,
            font=HEADER_FONT,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
        )

        dot_x = 0.31
        first_dot_y = list_y + collapsed_h + 0.07
        dot_ys = [first_dot_y + station_step * index for index in range(station_count)]
        if station_count > 1:
            station_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(dot_x - 0.014),
                Inches(dot_ys[0]),
                Inches(0.028),
                Inches(dot_ys[-1] - dot_ys[0]),
            )
            station_line.fill.solid()
            station_line.fill.fore_color.rgb = rgb(MUTED)
            station_line.line.fill.background()

        current_station = station_slides.index(current_slide)
        for station_index, dot_y in enumerate(dot_ys):
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(dot_x - 0.06),
                Inches(dot_y - 0.06),
                Inches(0.12),
                Inches(0.12),
            )
            if station_index < current_station:
                dot.fill.solid()
                dot.fill.fore_color.rgb = rgb(MUTED)
                dot.line.color.rgb = rgb(MUTED)
            elif station_index == current_station:
                dot.fill.solid()
                dot.fill.fore_color.rgb = rgb(SECONDARY)
                dot.line.color.rgb = rgb(SECONDARY)
            else:
                dot.fill.solid()
                dot.fill.fore_color.rgb = rgb(DEEP)
                dot.line.color.rgb = rgb(PALE)
            dot.line.width = Pt(1.1)

        list_y += active_h

    return rail


def normalize_all_act_rails(prs):
    """Replace legacy top-pinned rails while preserving all non-rail slide content."""
    rail_right = Inches(1.43)
    for physical_slide_number, slide in enumerate(prs.slides, start=1):
        rail_shapes = [
            shape
            for shape in slide.shapes
            if shape.left < rail_right and shape.left + shape.width <= rail_right
        ]
        for shape in rail_shapes:
            shape._element.getparent().remove(shape._element)
        if physical_slide_number == 22:
            conceptual_slide_number = 23
        elif physical_slide_number == 23:
            conceptual_slide_number = 23.5
        else:
            conceptual_slide_number = physical_slide_number
        if conceptual_slide_number in ACT_IV_STATIONS:
            station_slides = ACT_IV_STATIONS
        elif conceptual_slide_number in ACT_V_STATIONS:
            station_slides = ACT_V_STATIONS
        else:
            station_slides = None
        add_inline_act_rail(
            slide, conceptual_slide_number, station_slides=station_slides
        )


def add_body_card(slide, number, text, x, y, w, h, dark=False, font_size=15.0):
    fill = DEEP if dark else WHITE
    border = "334155" if dark else PALE
    text_color = OFFWHITE if dark else PRIMARY
    card = add_rect(slide, x, y, w, h, fill, border, radius=True, line_width=1.2)
    add_text(
        slide,
        f"{number:02d}",
        x + 0.18,
        y + 0.15,
        0.42,
        0.24,
        10,
        SECONDARY,
        font=HEADER_FONT,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    mark = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x + 0.18),
        Inches(y + 0.48),
        Inches(0.34),
        Inches(0.035),
    )
    mark.fill.solid()
    mark.fill.fore_color.rgb = rgb(SECONDARY)
    mark.line.fill.background()
    add_text(
        slide,
        text,
        x + 0.67,
        y + 0.14,
        w - 0.85,
        h - 0.25,
        font_size,
        text_color,
        font=BODY_FONT,
        valign=MSO_ANCHOR.MIDDLE,
    )
    return card


def add_visual_placeholder(slide, instruction, x, y, w, h, dark=False, split=False):
    fill = "182338" if dark else "E8EDF3"
    border = "516177" if dark else MUTED
    text_color = PALE if dark else PRIMARY
    add_rect(slide, x, y, w, h, fill, border, radius=True, line_width=1.2)

    if split:
        gap = 0.12
        half = (w - gap - 0.42) / 2
        left_x = x + 0.15
        right_x = left_x + half + gap
        inner_y = y + 0.52
        inner_h = h - 1.43
        add_rect(
            slide,
            left_x,
            inner_y,
            half,
            inner_h,
            "25334A" if dark else "F8FAFC",
            "516177" if dark else PALE,
            radius=True,
        )
        add_rect(
            slide,
            right_x,
            inner_y,
            half,
            inner_h,
            "131C2D" if dark else "CBD5E1",
            "516177" if dark else PALE,
            radius=True,
        )
        add_text(
            slide,
            "SHARP FRAME",
            left_x,
            inner_y + inner_h / 2 - 0.18,
            half,
            0.36,
            14,
            SECONDARY,
            font=HEADER_FONT,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            "WORST-BLURRED FRAME",
            right_x,
            inner_y + inner_h / 2 - 0.18,
            half,
            0.36,
            13,
            ACCENT,
            font=HEADER_FONT,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            "VISUAL PLACEHOLDER",
            x + 0.18,
            y + 0.13,
            w - 0.36,
            0.28,
            11,
            SECONDARY,
            font=HEADER_FONT,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            instruction,
            x + 0.20,
            y + h - 0.72,
            w - 0.40,
            0.56,
            10.5,
            text_color,
            font=BODY_FONT,
            italic=True,
            valign=MSO_ANCHOR.MIDDLE,
        )
    else:
        add_text(
            slide,
            "VISUAL PLACEHOLDER",
            x + 0.25,
            y + 0.28,
            w - 0.50,
            0.35,
            12,
            SECONDARY,
            font=HEADER_FONT,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            instruction,
            x + 0.35,
            y + 0.92,
            w - 0.70,
            h - 1.35,
            14,
            text_color,
            font=BODY_FONT,
            italic=True,
            valign=MSO_ANCHOR.MIDDLE,
        )
    return


def add_speech_bubble(slide, text, x, y, w, h, font_size=15.5):
    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    callout.name = "SpeechBubble_Callout"
    callout.fill.solid()
    callout.fill.fore_color.rgb = rgb(ACCENT)
    callout.line.color.rgb = rgb(WHITE)
    callout.line.width = Pt(1.2)
    for index, value in enumerate((0.45, 0.85, 0.16667)):
        callout.adjustments[index] = value

    frame = callout.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Inches(0.14)
    frame.margin_right = Inches(0.14)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    run = paragraph.add_run()
    run.text = text
    run.font.name = HEADER_FONT
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = rgb(WHITE)
    return callout


def _shape_fill_rgb(shape):
    try:
        return str(shape.fill.fore_color.rgb)
    except (AttributeError, TypeError, ValueError):
        return None


def prepare_speech_overlays(prs, mascot_path, serious_mascot_path):
    """Name, position, and populate the five animated speech overlays."""
    if not mascot_path.exists():
        raise FileNotFoundError(f"Missing mascot asset: {mascot_path}")
    if not serious_mascot_path.exists():
        raise FileNotFoundError(
            f"Missing serious mascot asset: {serious_mascot_path}"
        )

    for slide_number in SPEECH_OVERLAY_SLIDES:
        if slide_number > len(prs.slides):
            continue
        slide = prs.slides[slide_number - 1]
        if any(shape.name == "SpeechOverlay_Group" for shape in slide.shapes):
            continue

        native_callout = next(
            (shape for shape in slide.shapes if shape.name == "SpeechBubble_Callout"),
            None,
        )
        if native_callout is not None:
            if not any(shape.name == "Mascot" for shape in slide.shapes):
                slide_mascot_path = (
                    serious_mascot_path if slide_number == 11 else mascot_path
                )
                mascot = slide.shapes.add_picture(
                    str(slide_mascot_path),
                    Inches(MASCOT_X),
                    Inches(MASCOT_Y),
                    Inches(MASCOT_SIZE),
                    Inches(MASCOT_SIZE),
                )
                mascot.name = "Mascot"
            continue

        for shape in list(slide.shapes):
            if shape.name in ("Mascot", "Mascot_Origin"):
                shape._element.getparent().remove(shape._element)

        marker = SPEECH_TEXT_MARKERS[slide_number]
        bubble_text = next(
            (
                shape
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
                and marker in shape.text
                and shape.top >= Inches(5.4)
            ),
            None,
        )
        if bubble_text is None:
            raise ValueError(f"Could not find the speech text on slide {slide_number}.")

        accent_shapes = [
            shape
            for shape in slide.shapes
            if _shape_fill_rgb(shape) == ACCENT and shape.top >= Inches(5.4)
        ]
        bubble = next(
            (
                shape
                for shape in accent_shapes
                if shape.name == "SpeechBubble_Background"
                or shape.width >= Inches(2.0)
            ),
            None,
        )
        tail = next(
            (
                shape
                for shape in accent_shapes
                if shape.name == "SpeechBubble_Tail"
                or shape.width < Inches(1.0)
            ),
            None,
        )
        if bubble is None or tail is None:
            raise ValueError(f"Could not find all speech shapes on slide {slide_number}.")

        bubble_x = SPEECH_BUBBLE_RIGHT - bubble.width / Inches(1)
        bubble_w = bubble.width / Inches(1)
        bubble_h = bubble.height / Inches(1)
        bubble_y = SPEECH_TAIL_TARGET_Y - 0.28 - bubble_h + 0.04

        bubble.left = Inches(bubble_x)
        bubble.top = Inches(bubble_y)
        bubble.name = "SpeechBubble_Background"
        bubble_text.left = Inches(bubble_x + 0.22)
        bubble_text.top = Inches(bubble_y + 0.10)
        bubble_text.width = Inches(bubble_w - 0.44)
        bubble_text.height = Inches(bubble_h - 0.20)
        bubble_text.name = "SpeechBubble_Text"

        tail.left = Inches(SPEECH_BUBBLE_RIGHT - 0.22)
        tail.top = Inches(bubble_y + bubble_h - 0.04)
        tail.width = Inches(0.38)
        tail.height = Inches(0.28)
        tail.rotation = 180
        tail.name = "SpeechBubble_Tail"

        slide_mascot_path = (
            serious_mascot_path if slide_number == 11 else mascot_path
        )
        mascot = slide.shapes.add_picture(
            str(slide_mascot_path),
            Inches(MASCOT_X),
            Inches(MASCOT_Y),
            Inches(MASCOT_SIZE),
            Inches(MASCOT_SIZE),
        )
        mascot.name = "Mascot"


def bring_speech_overlays_to_front(prs):
    """Keep animated mascot/callout groups above subsequently inserted visuals."""
    for slide_number in SPEECH_OVERLAY_SLIDES:
        if slide_number > len(prs.slides):
            continue
        slide = prs.slides[slide_number - 1]
        overlay = next(
            (shape for shape in slide.shapes if shape.name == "SpeechOverlay_Group"),
            None,
        )
        if overlay is None:
            continue
        shape_tree = overlay._element.getparent()
        shape_tree.remove(overlay._element)
        shape_tree.append(overlay._element)

def apply_speech_animations(presentation_path, serious_mascot_path):
    """Replace speech parts with native callouts and add Zoom entrances."""
    if not ANIMATION_SCRIPT.exists():
        raise FileNotFoundError(f"Missing animation script: {ANIMATION_SCRIPT}")
    result = subprocess.run(
        [
            "powershell.exe",
            "-NoProfile",
            "-ExecutionPolicy",
            "Bypass",
            "-File",
            str(ANIMATION_SCRIPT),
            "-PresentationPath",
            str(presentation_path.resolve()),
            "-SeriousMascotPath",
            str(serious_mascot_path.resolve()),
        ],
        cwd=SCRIPT_DIR,
        check=True,
        capture_output=True,
        text=True,
    )
    if result.stdout.strip():
        print(result.stdout.strip())


def configure_slide17_video(presentation_path):
    """Make Slide 17's media-play effect automatic in the PowerPoint timeline."""
    presentation_path = Path(presentation_path).resolve()
    temporary_path = presentation_path.with_suffix(".autoplay.tmp.pptx")
    slide_part = "ppt/slides/slide17.xml"
    namespaces = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main"
    }

    with ZipFile(presentation_path, "r") as source:
        slide_xml = etree.fromstring(source.read(slide_part))
        play_commands = slide_xml.xpath(
            ".//p:cmd[starts-with(@cmd, 'playFrom')]", namespaces=namespaces
        )
        if len(play_commands) != 1:
            raise ValueError(
                f"Expected one Slide 17 media-play command, found {len(play_commands)}."
            )

        play_timing_node = play_commands[0].xpath(
            "ancestor::p:cTn[@nodeType='clickEffect'][1]", namespaces=namespaces
        )
        if len(play_timing_node) != 1:
            raise ValueError("Could not locate Slide 17's click-triggered media effect.")
        play_timing_node[0].set("nodeType", "withEffect")

        start_conditions = slide_xml.xpath(
            ".//p:cTn[@nodeType='mainSeq']/p:childTnLst/p:par/p:cTn/"
            "p:stCondLst/p:cond[@delay='indefinite']",
            namespaces=namespaces,
        )
        if len(start_conditions) != 1:
            raise ValueError(
                "Could not uniquely locate Slide 17's media start condition."
            )
        start_conditions[0].set("delay", "0")
        updated_slide_xml = etree.tostring(
            slide_xml, xml_declaration=True, encoding="UTF-8", standalone=True
        )

        with ZipFile(temporary_path, "w", compression=ZIP_DEFLATED) as target:
            for item in source.infolist():
                payload = (
                    updated_slide_xml if item.filename == slide_part else source.read(item)
                )
                target.writestr(item, payload)

    os.replace(temporary_path, presentation_path)
    print("Slide 17 video set to start automatically with the slide.")


def add_notes(slide, notes, source=None, visual_spec=None):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    if text_frame is None:
        # Decks that have passed through desktop PowerPoint can expose a notes
        # master whose placeholders are not cloned onto newly appended slides.
        # Clone the standard notes placeholders from an existing slide so the
        # new slide receives a proper notes body without changing prior notes.
        presentation = slide.part.package.presentation_part.presentation
        template_notes = next(
            (
                existing_slide.notes_slide
                for existing_slide in presentation.slides
                if existing_slide is not slide
                and existing_slide.notes_slide.notes_text_frame is not None
            ),
            None,
        )
        if template_notes is None:
            raise ValueError("No existing notes placeholder is available to clone.")
        for placeholder in template_notes.placeholders:
            if placeholder.placeholder_format.type in (
                PP_PLACEHOLDER.SLIDE_IMAGE,
                PP_PLACEHOLDER.BODY,
                PP_PLACEHOLDER.SLIDE_NUMBER,
            ):
                notes_slide.shapes.clone_placeholder(placeholder)
        text_frame = notes_slide.notes_text_frame
        if text_frame is None:
            raise ValueError("Could not create the notes text placeholder.")
    text_frame.text = notes
    if visual_spec:
        paragraph = text_frame.add_paragraph()
        paragraph.text = f"Visual: {visual_spec}"
    if source:
        paragraph = text_frame.add_paragraph()
        paragraph.text = f"[Sources]\n- {source}"


def add_picture_cover(slide, image_path, x, y, w, h):
    """Add an image using a centered 'cover' crop, like CSS object-fit: cover."""
    from PIL import Image

    with Image.open(image_path) as image:
        image_w, image_h = image.size
    image_ratio = image_w / image_h
    frame_ratio = w / h

    picture = slide.shapes.add_picture(
        str(image_path), Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if image_ratio > frame_ratio:
        visible_fraction = frame_ratio / image_ratio
        crop = (1.0 - visible_fraction) / 2.0
        picture.crop_left = crop
        picture.crop_right = crop
    elif image_ratio < frame_ratio:
        visible_fraction = image_ratio / frame_ratio
        crop = (1.0 - visible_fraction) / 2.0
        picture.crop_top = crop
        picture.crop_bottom = crop
    return picture


def add_rounded_picture_cover(
    slide, image_path, x, y, w, h, line_color="516177", line_width=1.5
):
    """Add a cover-cropped image with rounded corners and a visible frame."""
    picture = add_picture_cover(slide, image_path, x, y, w, h)
    picture._element.spPr.prstGeom.set("prst", "roundRect")

    frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    frame.fill.background()
    frame.line.color.rgb = rgb(line_color)
    frame.line.width = Pt(line_width)
    return picture


def add_picture_contain(slide, image_path, x, y, w, h, dark=False):
    """Fit a complete image inside a rounded presentation frame."""
    from PIL import Image

    frame_fill = "182338" if dark else "E8EDF3"
    frame_line = "516177" if dark else MUTED
    add_rect(slide, x, y, w, h, frame_fill, frame_line, radius=True, line_width=1.2)

    with Image.open(image_path) as image:
        image_w, image_h = image.size
    scale = min((w - 0.20) / image_w, (h - 0.20) / image_h)
    placed_w = image_w * scale
    placed_h = image_h * scale
    placed_x = x + (w - placed_w) / 2
    placed_y = y + (h - placed_h) / 2
    picture = slide.shapes.add_picture(
        str(image_path),
        Inches(placed_x),
        Inches(placed_y),
        Inches(placed_w),
        Inches(placed_h),
    )
    picture._element.spPr.prstGeom.set("prst", "roundRect")
    return picture


def add_movie_contain(
    slide, movie_path, poster_path, x, y, w, h, shape_name, dark=False
):
    """Embed a movie without distorting its source aspect ratio."""
    from PIL import Image

    frame_fill = "182338" if dark else "E8EDF3"
    frame_line = "516177" if dark else MUTED
    add_rect(slide, x, y, w, h, frame_fill, frame_line, radius=True, line_width=1.2)

    with Image.open(poster_path) as image:
        image_w, image_h = image.size
    scale = min((w - 0.12) / image_w, (h - 0.12) / image_h)
    placed_w = image_w * scale
    placed_h = image_h * scale
    placed_x = x + (w - placed_w) / 2
    placed_y = y + (h - placed_h) / 2
    movie = slide.shapes.add_movie(
        str(movie_path),
        Inches(placed_x),
        Inches(placed_y),
        Inches(placed_w),
        Inches(placed_h),
        poster_frame_image=str(poster_path),
        mime_type="video/mp4",
    )
    movie.name = shape_name
    return movie


def build_slide_1(prs, visual_1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        str(visual_1), 0, 0, width=prs.slide_width, height=prs.slide_height
    )

    add_rect(slide, 1.58, 0.32, 8.36, 1.68, DEEP, None, radius=True)
    add_text(
        slide,
        "The Raycast Challenge",
        1.84,
        0.53,
        7.80,
        0.65,
        50,
        WHITE,
        font=HEADER_FONT,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "Pick a pixel in one drone frame. Find the same point in twelve others. ≤10 pixels or it doesn't count.",
        1.86,
        1.27,
        7.82,
        0.48,
        19.5,
        OFFWHITE,
        font=BODY_FONT,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_act_rail(slide, 1)
    add_speech_bubble(slide, '"≤10 pixels in your dreams, mate."', 9.27, 6.20, 3.55, 0.82)
    add_notes(
        slide,
        'Open cold with the image, not with the org chart. This is a "here\'s the trick" opener — the rest of the deck is "here\'s what it took."',
        "User-provided visual: Raycast_Slide_1_Visual.png",
        "A single dramatic frame from the footage — the van, mid-field, fisheye-distorted — with a crosshair on it and thin lines fanning out to ~12 hovering drones (cartoon). This is the whole pitch in one image. [NTPvChat]",
    )


def build_slide_2(prs, visual_2):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, OFFWHITE)
    add_title(slide, "The Assignment", dark=False)

    add_picture_cover(slide, visual_2, 1.75, 1.42, 4.18, 5.55)

    items = [
        "A sequence of drone FPV frames, one stationary van, one agricultural field",
        "Several frames badly motion/defocus-blurred",
        "Deliverable: source code, a methodology writeup, and proof — not just a claim — of ≤10px reprojection accuracy",
        'Complete architectural freedom. "Ingenuity and a tinkerer mindset," officially encouraged.',
    ]
    positions = [
        (6.22, 1.42),
        (9.57, 1.42),
        (6.22, 4.28),
        (9.57, 4.28),
    ]
    for index, (text, (x, y)) in enumerate(zip(items, positions), start=1):
        add_body_card(slide, index, text, x, y, 3.05, 2.69, dark=False)

    add_act_rail(slide, 2)
    add_notes(
        slide,
        "This slide exists so the audience knows the constraints going in: freedom, blur, repetition, and a hard numeric bar.",
    )


def build_slide_3(prs, sharp_frame, blurred_frame):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PRIMARY)
    add_title(slide, "The Challenges", dark=True)

    visual_instruction = (
        "Side-by-side crop of a sharp frame vs. one of the worst-blurred frames, same rough patch of ground. [NTPvChat]"
    )
    add_rect(slide, 1.75, 1.42, 6.15, 5.55, "182338", "516177", radius=True, line_width=1.2)
    add_rounded_picture_cover(slide, sharp_frame, 1.90, 1.94, 2.805, 4.12)
    add_rounded_picture_cover(slide, blurred_frame, 4.825, 1.94, 2.805, 4.12)
    add_rect(slide, 1.90, 1.94, 2.805, 0.48, PRIMARY, None, radius=True)
    add_text(
        slide, "SHARP FRAME", 1.90, 1.94, 2.805, 0.48, 14, SECONDARY,
        font=HEADER_FONT, bold=True, align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_rect(slide, 4.825, 1.94, 2.805, 0.48, PRIMARY, None, radius=True)
    add_text(
        slide, "WORST-BLURRED FRAME", 4.825, 1.94, 2.805, 0.48, 13, ACCENT,
        font=HEADER_FONT, bold=True, align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide, visual_instruction, 1.95, 6.25, 5.75, 0.56, 10.5, PALE,
        font=BODY_FONT, italic=True, valign=MSO_ANCHOR.MIDDLE,
    )

    items = [
        "The field is highly repetitive — one patch of dirt looks like every other patch of dirt",
        "Several frames are so blurred that even a human struggles to place a feature confidently",
        "GPS telemetry on the drones is off by multiple metres — you can't just trust the numbers you're handed",
    ]
    ys = [1.42, 3.28, 5.14]
    for index, (text, y) in enumerate(zip(items, ys), start=1):
        add_body_card(slide, index, text, 8.22, y, 4.60, 1.65, dark=True)

    add_act_rail(slide, 3)
    add_notes(
        slide,
        "This is the thesis of the whole deck: every subsequent architectural decision is a response to one of these three problems.",
    )


def build_slide_4(prs, visual_4):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PRIMARY)
    add_title(
        slide,
        "First Commit: Already a Real Pipeline (Thanks, Claude :-))",
        dark=True,
        font_size=31,
    )

    visual_spec = (
        "Load Frames → OCR (reads the HUD overlay) → Undistort → Pose "
        "(GPS→ENU + prior rotation) → Refine (`pitch_optimizer.py`) → "
        "Interactive Viewer (ray-cast on click).\nRefine → Blender. [NTPvChat]"
    )
    add_picture_contain(slide, visual_4, 1.72, 1.42, 2.85, 5.42, dark=True)

    items = [
        "OCR reads GPS/heading/altitude straight off the HUD overlay",
        "GPS → local ENU coordinates",
        "Fisheye lens undistortion",
        "A hand-rolled pitch optimizer",
        "Ray–ground intersection math",
        "An interactive click-to-reproject viewer",
        "Debug argument: Export solve to Blender",
    ]
    positions = [
        (4.82, 1.48),
        (8.78, 1.48),
        (4.82, 2.65),
        (8.78, 2.65),
        (4.82, 3.82),
        (8.78, 3.82),
        (4.82, 4.99),
    ]
    for index, (text, (x, y)) in enumerate(zip(items, positions), start=1):
        add_body_card(
            slide, index, text, x, y, 3.66, 1.02, dark=True, font_size=12.7
        )

    add_act_i_rail(slide, 4)
    add_speech_bubble(
        slide,
        '"First commit after ~4 workday... literally a war-crime."',
        8.38,
        6.24,
        4.36,
        0.72,
    )
    add_notes(
        slide,
        "Establish early that day one wasn't a toy — it was a mature skeleton that mostly survived. Mention roll bug which was easily spotted in Blender.",
        visual_spec=visual_spec,
    )


def build_slide_5(prs, visual_5):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, OFFWHITE)
    add_title(slide, "Lesson #1, Learned in the First Hour", font_size=36)

    body = (
        'At this stage, pitch is computed from the ground plane — and the van is '
        'a real 3D object sitting on top of that plane, not part of it. Mask it '
        'out, or "what does flat ground look like" gets answered by the wrong geometry.'
    )
    add_rect(slide, 1.75, 1.48, 4.22, 4.88, WHITE, PALE, radius=True, line_width=1.2)
    add_text(
        slide,
        body,
        2.05,
        1.84,
        3.62,
        4.16,
        20,
        PRIMARY,
        font=BODY_FONT,
        valign=MSO_ANCHOR.MIDDLE,
    )

    visual_spec = (
        'A frame with the van bounding-boxed and crossed out in red, labeled '
        '"real geometry — wrong category." [NTPvChat]'
    )
    add_rounded_picture_cover(
        slide, visual_5, 6.23, 1.48, 6.57, 4.65, line_color=MUTED, line_width=1.2
    )

    add_act_i_rail(slide, 5)
    add_speech_bubble(
        slide,
        '"You have no idea how important those masked-out van pixels will become towards the middle of the project and onwards."',
        6.98,
        6.25,
        5.76,
        0.82,
    )
    add_notes(
        slide,
        "This is specifically about excluding a real foreground object from a ground-only computation — don't conflate it with the screen-overlay masking that shows up later (Slide 7 onward); that's a related-but-distinct lesson with its own payoff.",
        visual_spec=visual_spec,
    )


def replace_slide_5_placeholder(prs, visual_5):
    """Replace the existing Slide 5 placeholder without rebuilding prior slides."""
    slide = prs.slides[4]
    visual_left = Inches(6.20)
    visual_top = Inches(1.40)
    visual_bottom = Inches(6.10)

    existing_picture = any(
        shape.shape_type == 13
        and shape.left >= visual_left
        and shape.top >= visual_top
        and shape.top < visual_bottom
        for shape in slide.shapes
    )
    if existing_picture:
        return

    placeholder_shapes = [
        shape
        for shape in slide.shapes
        if shape.left >= visual_left
        and shape.top >= visual_top
        and shape.top < visual_bottom
    ]
    for shape in placeholder_shapes:
        shape._element.getparent().remove(shape._element)

    add_rounded_picture_cover(
        slide, visual_5, 6.23, 1.48, 6.57, 4.65, line_color=MUTED, line_width=1.2
    )


def replace_slide_visual(
    prs, slide_number, image_path, x, y, w, h, shape_name, dark=True
):
    """Replace a named slide's visual placeholder without disturbing its text."""
    slide = prs.slides[slide_number - 1]
    if any(shape.name == shape_name for shape in slide.shapes):
        return

    left = Inches(x - 0.04)
    top = Inches(y - 0.04)
    right = Inches(x + w + 0.04)
    bottom = Inches(y + h + 0.04)
    placeholder_shapes = [
        shape
        for shape in slide.shapes
        if shape.left >= left
        and shape.top >= top
        and shape.left + shape.width <= right
        and shape.top + shape.height <= bottom
    ]
    for shape in placeholder_shapes:
        shape._element.getparent().remove(shape._element)

    picture = add_picture_contain(slide, image_path, x, y, w, h, dark=dark)
    picture.name = shape_name


def replace_slide_video(
    prs,
    slide_number,
    movie_path,
    poster_path,
    x,
    y,
    w,
    h,
    shape_name,
    dark=False,
):
    """Replace a slide visual placeholder with an embedded movie."""
    slide = prs.slides[slide_number - 1]
    if any(shape.name == shape_name for shape in slide.shapes):
        return

    left = Inches(x - 0.04)
    top = Inches(y - 0.04)
    right = Inches(x + w + 0.04)
    bottom = Inches(y + h + 0.04)
    placeholder_shapes = [
        shape
        for shape in slide.shapes
        if shape.left >= left
        and shape.top >= top
        and shape.left + shape.width <= right
        and shape.top + shape.height <= bottom
    ]
    for shape in placeholder_shapes:
        shape._element.getparent().remove(shape._element)

    add_movie_contain(
        slide,
        movie_path,
        poster_path,
        x,
        y,
        w,
        h,
        shape_name,
        dark=dark,
    )


def build_slide_6(prs, visual_6):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PRIMARY)
    add_title(
        slide,
        '"Lots of Improvements Tried. Result Still Not Amazing."',
        dark=True,
        font_size=28,
    )

    add_rect(
        slide, 1.75, 1.48, 4.48, 4.98, DEEP, "334155", radius=True, line_width=1.2
    )
    add_rich_text(
        slide,
        [
            ("Day two: the solver gets switched to Ceres (", BODY_FONT),
            ("pyceres", "Cascadia Mono"),
            (
                ") — real nonlinear least-squares, not a hand-rolled optimizer. "
                "Two tools are born here that outlive nearly everything else in "
                "the codebase: ",
                BODY_FONT,
            ),
            ("camera_deltas.py", "Cascadia Mono"),
            (
                " (solved-vs-telemetry sanity check) and a GeoCalib-based "
                "pitch/roll estimator.",
                BODY_FONT,
            ),
        ],
        2.06,
        1.82,
        3.86,
        4.30,
        19,
        OFFWHITE,
        valign=MSO_ANCHOR.MIDDLE,
    )

    visual_spec = "User-provided visual: Raycast_Slide_6_Visual.png"
    picture = add_picture_contain(
        slide, visual_6, 6.50, 1.48, 6.30, 4.98, dark=True
    )
    picture.name = "Slide6_Visual"

    add_act_i_rail(slide, 6)
    add_notes(
        slide,
        'This is the first "one twin lives, one twin fades" beat — good rhythm-setter for the whole deck\'s structure.',
        visual_spec=visual_spec,
    )


def build_slide_7(prs, visual_7):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PRIMARY)
    add_title(slide, "Ten Manual Graphics Overlay Masks", dark=True, font_size=36)

    add_rect(
        slide, 1.75, 1.48, 4.45, 4.98, DEEP, "334155", radius=True, line_width=1.2
    )
    add_rich_text(
        slide,
        [
            ("HUD_REGIONS", "Cascadia Mono"),
            (
                ": 10 rectangles, geometrically masking every HUD overlay element "
                "in raw camera pixels. The top-centre navigation bar alone is split "
                "into three separate boxes — because one wide box, dragged through "
                "fisheye undistortion, would bow at the corners and expose exactly "
                "the content it was supposed to hide.",
                BODY_FONT,
            ),
        ],
        2.06,
        1.82,
        3.83,
        4.30,
        18,
        OFFWHITE,
        valign=MSO_ANCHOR.MIDDLE,
    )

    visual_spec = "User-provided visual: Raycast_Slide_7_Visual.png"
    picture = add_picture_contain(
        slide, visual_7, 6.48, 1.48, 6.32, 4.98, dark=True
    )
    picture.name = "Slide7_Visual"
    add_act_ii_rail(slide, 7)
    add_notes(slide, "", visual_spec=visual_spec)


def add_video_placeholder(slide, x, y, w, h):
    """Add a clean 16:9 frame ready for a future embedded video."""
    frame = add_rect(
        slide, x, y, w, h, "E8EDF3", MUTED, radius=True, line_width=1.2
    )
    frame.name = "Slide8_Video_Placeholder"

    circle_size = 0.78
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x + (w - circle_size) / 2),
        Inches(y + (h - circle_size) / 2 - 0.08),
        Inches(circle_size),
        Inches(circle_size),
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = rgb(PRIMARY)
    circle.line.fill.background()

    play = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(x + w / 2 - 0.13),
        Inches(y + h / 2 - 0.23),
        Inches(0.28),
        Inches(0.34),
    )
    play.rotation = 90
    play.fill.solid()
    play.fill.fore_color.rgb = rgb(WHITE)
    play.line.fill.background()
    add_text(
        slide,
        "LANDSCAPE VIDEO PLACEHOLDER",
        x + 0.35,
        y + h - 0.50,
        w - 0.70,
        0.26,
        11,
        PRIMARY,
        font=HEADER_FONT,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_slide_8_stage_card(slide, number, text, x, font_size):
    """Add one compact stage in Slide 8's horizontal process row."""
    y, w, h = 5.72, 2.68, 1.42
    add_rect(slide, x, y, w, h, WHITE, PALE, radius=True, line_width=1.0)
    add_text(
        slide,
        f"{number:02d}",
        x + 0.15,
        y + 0.10,
        0.36,
        0.20,
        9.2,
        SECONDARY,
        font=HEADER_FONT,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    mark = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x + 0.15),
        Inches(y + 0.35),
        Inches(0.34),
        Inches(0.03),
    )
    mark.fill.solid()
    mark.fill.fore_color.rgb = rgb(SECONDARY)
    mark.line.fill.background()
    add_text(
        slide,
        text,
        x + 0.16,
        y + 0.44,
        w - 0.32,
        h - 0.53,
        font_size,
        PRIMARY,
        font=BODY_FONT,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0.02,
    )


def layout_slide_8(slide, video_path=None, poster_path=None):
    """Lay out Slide 8 around a central landscape video."""
    add_title(slide, "The Manual Correspondence Picker Is Born", font_size=35)
    if video_path is not None and poster_path is not None:
        movie = slide.shapes.add_movie(
            str(video_path),
            Inches(3.70),
            Inches(1.35),
            Inches(7.25),
            Inches(4.08),
            poster_frame_image=str(poster_path),
            mime_type="video/mp4",
        )
        movie.name = "Slide8_Video"
    else:
        add_video_placeholder(slide, 3.70, 1.35, 7.25, 4.08)

    items = [
        "Multi-frame grid, independent zoom/pan, adjustable marker size",
        (
            "Four feature types, cycled with a key: ground, roof, wheel_axis, "
            "roof_edge — the last two are pairs with known real-world distances "
            "(wheelbase, van width) baked in as hard constraints"
        ),
        (
            "Right-click removes a single point; d deletes a whole correspondence "
            "or just one frame's mark, depending on context; [/] step through saved "
            "correspondences filtered by type"
        ),
        (
            "A quit-safety net: unsaved changes require a second keypress to "
            "actually discard"
        ),
    ]
    positions = [1.75, 4.54, 7.33, 10.12]
    font_sizes = [10.0, 8.8, 8.6, 9.4]
    for index, (text, x, font_size) in enumerate(
        zip(items, positions, font_sizes), start=1
    ):
        add_slide_8_stage_card(slide, index, text, x, font_size)


def normalize_slide_8_layout(prs, video_path=None, poster_path=None):
    """Rebuild Slide 8's content area without disturbing its rail or notes."""
    if len(prs.slides) < 8:
        return
    slide = prs.slides[7]
    content_shapes = [
        shape for shape in slide.shapes if shape.left >= Inches(1.43)
    ]
    for shape in content_shapes:
        shape._element.getparent().remove(shape._element)
    set_background(slide, OFFWHITE)
    layout_slide_8(slide, video_path, poster_path)


def configure_video_click_timing(
    prs, slide_number, shape_name, duration_ms, followup_zoom_shape_name=None
):
    """Play an embedded video on the next click/spacebar press, full-screen."""
    if len(prs.slides) < slide_number:
        return
    slide = prs.slides[slide_number - 1]
    video = next((shape for shape in slide.shapes if shape.name == shape_name), None)
    if video is None:
        return

    slide_element = slide._element
    existing_timing = slide_element.find(qn("p:timing"))
    if existing_timing is not None:
        slide_element.remove(existing_timing)

    shape_id = video.shape_id
    followup_xml = ""
    media_node_id = 7
    if followup_zoom_shape_name is not None:
        followup = next(
            (
                shape
                for shape in slide.shapes
                if shape.name == followup_zoom_shape_name
            ),
            None,
        )
        if followup is not None:
            followup_id = followup.shape_id
            media_node_id = 13
            followup_xml = f"""
                        <p:par>
                          <p:cTn id="7" fill="hold">
                            <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
                            <p:childTnLst>
                              <p:par>
                                <p:cTn id="8" fill="hold">
                                  <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                  <p:childTnLst>
                                    <p:par>
                                      <p:cTn id="9" presetID="23" presetClass="entr" presetSubtype="16" accel="12000" decel="18000" fill="hold" nodeType="clickEffect">
                                        <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                        <p:childTnLst>
                                          <p:set>
                                            <p:cBhvr>
                                              <p:cTn id="10" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
                                              <p:tgtEl><p:spTgt spid="{followup_id}"/></p:tgtEl>
                                              <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                                            </p:cBhvr>
                                            <p:to><p:strVal val="visible"/></p:to>
                                          </p:set>
                                          <p:anim calcmode="lin" valueType="num">
                                            <p:cBhvr>
                                              <p:cTn id="11" dur="250" fill="hold"/>
                                              <p:tgtEl><p:spTgt spid="{followup_id}"/></p:tgtEl>
                                              <p:attrNameLst><p:attrName>ppt_w</p:attrName></p:attrNameLst>
                                            </p:cBhvr>
                                            <p:tavLst><p:tav tm="0"><p:val><p:fltVal val="0"/></p:val></p:tav><p:tav tm="100000"><p:val><p:strVal val="#ppt_w"/></p:val></p:tav></p:tavLst>
                                          </p:anim>
                                          <p:anim calcmode="lin" valueType="num">
                                            <p:cBhvr>
                                              <p:cTn id="12" dur="250" fill="hold"/>
                                              <p:tgtEl><p:spTgt spid="{followup_id}"/></p:tgtEl>
                                              <p:attrNameLst><p:attrName>ppt_h</p:attrName></p:attrNameLst>
                                            </p:cBhvr>
                                            <p:tavLst><p:tav tm="0"><p:val><p:fltVal val="0"/></p:val></p:tav><p:tav tm="100000"><p:val><p:strVal val="#ppt_h"/></p:val></p:tav></p:tavLst>
                                          </p:anim>
                                        </p:childTnLst>
                                      </p:cTn>
                                    </p:par>
                                  </p:childTnLst>
                                </p:cTn>
                              </p:par>
                            </p:childTnLst>
                          </p:cTn>
                        </p:par>
            """
    timing = parse_xml(
        f"""
        <p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
          <p:tnLst>
            <p:par>
              <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
                <p:childTnLst>
                  <p:seq concurrent="1" nextAc="seek">
                    <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                      <p:childTnLst>
                        <p:par>
                          <p:cTn id="3" fill="hold">
                            <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
                            <p:childTnLst>
                              <p:par>
                                <p:cTn id="4" fill="hold">
                                  <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                  <p:childTnLst>
                                    <p:par>
                                      <p:cTn id="5" presetID="1" presetClass="mediacall" presetSubtype="0" fill="hold" nodeType="clickEffect">
                                        <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                        <p:childTnLst>
                                          <p:cmd type="call" cmd="playFrom(0.0)">
                                            <p:cBhvr>
                                              <p:cTn id="6" dur="{duration_ms}" fill="hold"/>
                                              <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                                            </p:cBhvr>
                                          </p:cmd>
                                        </p:childTnLst>
                                      </p:cTn>
                                    </p:par>
                                  </p:childTnLst>
                                </p:cTn>
                              </p:par>
                            </p:childTnLst>
                          </p:cTn>
                        </p:par>
                        {followup_xml}
                      </p:childTnLst>
                    </p:cTn>
                    <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
                    <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
                  </p:seq>
                  <p:video fullScrn="true">
                    <p:cMediaNode vol="80000">
                      <p:cTn id="{media_node_id}" fill="hold" display="0">
                        <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
                      </p:cTn>
                      <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                    </p:cMediaNode>
                  </p:video>
                </p:childTnLst>
              </p:cTn>
            </p:par>
          </p:tnLst>
        </p:timing>
        """
    )
    slide_element.append(timing)


def configure_video_autoplay_loop_timing(prs, slide_number, shape_name):
    """Start a video with its slide and loop it indefinitely."""
    if len(prs.slides) < slide_number:
        return
    slide = prs.slides[slide_number - 1]
    video = next((shape for shape in slide.shapes if shape.name == shape_name), None)
    if video is None:
        raise ValueError(f"Could not find {shape_name!r} on Slide {slide_number}.")

    slide_element = slide._element
    existing_timing = slide_element.find(qn("p:timing"))
    if existing_timing is not None:
        slide_element.remove(existing_timing)
    shape_id = video.shape_id
    timing = parse_xml(
        f"""
        <p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
          <p:tnLst>
            <p:par>
              <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
                <p:childTnLst>
                  <p:seq concurrent="1" nextAc="seek">
                    <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                      <p:childTnLst>
                        <p:par>
                          <p:cTn id="3" fill="hold">
                            <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                            <p:childTnLst>
                              <p:par>
                                <p:cTn id="4" fill="hold">
                                  <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                  <p:childTnLst>
                                    <p:par>
                                      <p:cTn id="5" presetID="1" presetClass="mediacall" presetSubtype="0" fill="hold" nodeType="withEffect">
                                        <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                        <p:childTnLst>
                                          <p:cmd type="call" cmd="playFrom(0.0)">
                                            <p:cBhvr>
                                              <p:cTn id="6" dur="1" fill="hold"/>
                                              <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                                            </p:cBhvr>
                                          </p:cmd>
                                        </p:childTnLst>
                                      </p:cTn>
                                    </p:par>
                                  </p:childTnLst>
                                </p:cTn>
                              </p:par>
                            </p:childTnLst>
                          </p:cTn>
                        </p:par>
                      </p:childTnLst>
                    </p:cTn>
                    <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
                    <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
                  </p:seq>
                  <p:video>
                    <p:cMediaNode vol="80000">
                      <p:cTn id="7" repeatCount="indefinite" fill="remove" display="0">
                        <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
                      </p:cTn>
                      <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                    </p:cMediaNode>
                  </p:video>
                </p:childTnLst>
              </p:cTn>
            </p:par>
          </p:tnLst>
        </p:timing>
        """
    )
    slide_element.append(timing)


def externalize_linked_video(presentation_path, slide_number, video_path):
    """Replace one embedded movie payload with a local external file link."""
    from externalize_pptx_videos import externalize_slide_video

    presentation_path = Path(presentation_path).resolve()
    video_path = Path(video_path).resolve()
    temporary_path = presentation_path.with_suffix(".external-video.tmp.pptx")
    with ZipFile(presentation_path, "r") as source:
        infos = source.infolist()
        entries = {info.filename: source.read(info.filename) for info in infos}
    removed_media = externalize_slide_video(entries, slide_number, video_path)
    if not removed_media:
        raise ValueError(f"No embedded video found on Slide {slide_number}.")
    with ZipFile(temporary_path, "w", allowZip64=True) as target:
        for info in infos:
            if info.filename in removed_media:
                continue
            target.writestr(info, entries[info.filename])
    os.replace(temporary_path, presentation_path)
    print(
        f"Slide {slide_number} video linked externally to {video_path}."
    )


def group_speech_overlay_with_mascot_pivot(slide, visible_shapes):
    """Group a speech overlay around an invisible mascot-centered pivot frame."""
    visible_shapes = list(visible_shapes)
    mascot = next((shape for shape in visible_shapes if shape.name == "Mascot"), None)
    if mascot is None:
        raise ValueError("Speech overlay cannot be grouped without a Mascot shape.")

    mascot_origin_x = mascot.left + mascot.width // 2
    mascot_origin_y = mascot.top + mascot.height // 2
    visible_left = min(shape.left for shape in visible_shapes)
    visible_top = min(shape.top for shape in visible_shapes)
    visible_right = max(shape.left + shape.width for shape in visible_shapes)
    visible_bottom = max(shape.top + shape.height for shape in visible_shapes)
    pivot_half_width = max(
        mascot_origin_x - visible_left,
        visible_right - mascot_origin_x,
    )
    pivot_half_height = max(
        mascot_origin_y - visible_top,
        visible_bottom - mascot_origin_y,
    )
    pivot_frame = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        mascot_origin_x - pivot_half_width,
        mascot_origin_y - pivot_half_height,
        2 * pivot_half_width,
        2 * pivot_half_height,
    )
    pivot_frame.name = "Mascot_Pivot_Frame"
    pivot_frame.fill.background()
    pivot_frame.line.fill.background()

    overlay_group = slide.shapes.add_group_shape(
        [*visible_shapes, pivot_frame]
    )
    overlay_group.name = "SpeechOverlay_Group"
    return overlay_group


def configure_speech_overlay_click_timing(prs, slide_number):
    """Reveal all speech-overlay pieces together on one click without COM."""
    if len(prs.slides) < slide_number:
        return
    slide = prs.slides[slide_number - 1]
    overlay_group = next(
        (shape for shape in slide.shapes if shape.name == "SpeechOverlay_Group"),
        None,
    )
    if overlay_group is not None:
        overlay_shapes = [overlay_group]
    else:
        overlay_names = (
            "SpeechBubble_Callout",
            "SpeechBubble_Background",
            "SpeechBubble_Text",
            "SpeechBubble_Tail",
            "Mascot",
        )
        overlay_shapes = [
            next((shape for shape in slide.shapes if shape.name == name), None)
            for name in overlay_names
        ]
        overlay_shapes = [shape for shape in overlay_shapes if shape is not None]
        if len(overlay_shapes) > 1:
            overlay_group = group_speech_overlay_with_mascot_pivot(
                slide, overlay_shapes
            )
            overlay_shapes = [overlay_group]
    if not overlay_shapes:
        return

    child_nodes = []
    timing_id = 6
    for shape in overlay_shapes:
        shape_id = shape.shape_id
        child_nodes.append(
            f"""
            <p:set>
              <p:cBhvr>
                <p:cTn id="{timing_id}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
                <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
              </p:cBhvr>
              <p:to><p:strVal val="visible"/></p:to>
            </p:set>
            <p:anim calcmode="lin" valueType="num">
              <p:cBhvr>
                <p:cTn id="{timing_id + 1}" dur="250" fill="hold"/>
                <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                <p:attrNameLst><p:attrName>ppt_w</p:attrName></p:attrNameLst>
              </p:cBhvr>
              <p:tavLst><p:tav tm="0"><p:val><p:fltVal val="0"/></p:val></p:tav><p:tav tm="100000"><p:val><p:strVal val="#ppt_w"/></p:val></p:tav></p:tavLst>
            </p:anim>
            <p:anim calcmode="lin" valueType="num">
              <p:cBhvr>
                <p:cTn id="{timing_id + 2}" dur="250" fill="hold"/>
                <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                <p:attrNameLst><p:attrName>ppt_h</p:attrName></p:attrNameLst>
              </p:cBhvr>
              <p:tavLst><p:tav tm="0"><p:val><p:fltVal val="0"/></p:val></p:tav><p:tav tm="100000"><p:val><p:strVal val="#ppt_h"/></p:val></p:tav></p:tavLst>
            </p:anim>
            """
        )
        timing_id += 3

    slide_element = slide._element
    existing_timing = slide_element.find(qn("p:timing"))
    if existing_timing is not None:
        slide_element.remove(existing_timing)
    timing = parse_xml(
        f"""
        <p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
          <p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>
            <p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>
              <p:par><p:cTn id="3" fill="hold"><p:stCondLst><p:cond delay="indefinite"/></p:stCondLst><p:childTnLst>
                <p:par><p:cTn id="4" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>
                  <p:par><p:cTn id="5" presetID="23" presetClass="entr" presetSubtype="16" accel="12000" decel="18000" fill="hold" nodeType="clickEffect"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst>
                    {''.join(child_nodes)}
                  </p:childTnLst></p:cTn></p:par>
                </p:childTnLst></p:cTn></p:par>
              </p:childTnLst></p:cTn></p:par>
            </p:childTnLst></p:cTn>
            <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
            <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
            </p:seq>
          </p:childTnLst></p:cTn></p:par></p:tnLst>
        </p:timing>
        """
    )
    slide_element.append(timing)


def build_slide_8(prs, video_path=None, poster_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, OFFWHITE)
    layout_slide_8(slide, video_path, poster_path)
    add_act_ii_rail(slide, 8)
    add_notes(
        slide,
        "This is worth a real video, not a mockup — it genuinely looks like a "
        "purpose-built desktop app, which surprises people expecting a debug script.",
        visual_spec="Landscape video of the manual correspondence picker in action.",
    )


def build_slide_9(prs, visual_9=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PRIMARY)
    add_title(slide, "The Scorer: A Heads-Up, Not a Verdict", dark=True, font_size=36)

    body = (
        'For every correspondence, crop the clicked point in each frame, embed '
        'with CLIP, average pairwise similarity. A low score is genuinely '
        'ambiguous: it can mean "wrong pick" — or "right pick, just doesn\'t look '
        'alike from this angle/blur." CLIP only sees pixels, not the context that '
        "told you it's the same rivet."
    )
    add_rect(
        slide, 1.75, 1.48, 4.52, 4.98, DEEP, "334155", radius=True, line_width=1.2
    )
    add_text(
        slide,
        body,
        2.07,
        1.82,
        3.88,
        4.30,
        18.5,
        OFFWHITE,
        font=BODY_FONT,
        valign=MSO_ANCHOR.MIDDLE,
    )

    visual_spec = (
        'Three crop pairs side by side, unlabeled at first glance: (1) high score, '
        'obviously the same feature; (2) low score, still correct — just an '
        'oblique-angle or blur mismatch; (3) low score, actually a mis-click. '
        'The point: (2) and (3) look equally "bad" to the algorithm. Only a human '
        'looking at the full frame — not just the crop — can tell them apart. '
        "[NTPvCode]"
    )
    if visual_9 is not None:
        picture = add_picture_contain(
            slide, visual_9, 6.54, 1.48, 6.26, 4.98, dark=True
        )
        picture.name = "Slide9_Visual"
    else:
        add_visual_placeholder(
            slide, visual_spec, 6.54, 1.48, 6.26, 4.98, dark=True
        )
    add_act_ii_rail(slide, 9)
    add_notes(
        slide,
        "Emphasize what it's not: not a 3D geometric check, purely \"do these crops look alike\" — and that's a real limitation, not a footnote. The project's own README says as much: low scores across oblique frame pairs are expected, not necessarily wrong.",
        visual_spec=visual_spec,
    )


def build_slide_10(prs, visual_10=None, poster_10=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, OFFWHITE)
    add_title(
        slide, '"Manual Feature Matching Works Great on 7 Frames"', font_size=28
    )

    visual_spec = (
        "The best-looking proof-sheet reprojection from this era, one source pixel "
        "accurately landing across all 7 frames. [NTPvCode]"
    )
    if visual_10 is not None and poster_10 is not None:
        add_movie_contain(
            slide,
            visual_10,
            poster_10,
            1.75,
            1.48,
            6.18,
            4.82,
            "Slide10_Video",
            dark=False,
        )
    else:
        add_visual_placeholder(
            slide, visual_spec, 1.75, 1.48, 6.18, 4.82, dark=False
        )

    body = (
        "Known van geometry (wheel axis, roof edges) as hard anchors + hand-picked "
        "ground correspondences → a genuinely good calibration on the 7 clearest "
        "frames. This becomes the quality bar the rest of the project measures "
        "itself against."
    )
    add_rect(slide, 8.20, 1.48, 4.60, 4.82, WHITE, PALE, radius=True, line_width=1.2)
    add_text(
        slide,
        body,
        8.52,
        1.84,
        3.96,
        4.08,
        19,
        PRIMARY,
        font=BODY_FONT,
        valign=MSO_ANCHOR.MIDDLE,
    )

    add_act_ii_rail(slide, 10)
    add_speech_bubble(
        slide,
        '"This exact commit also complains, in its own message, about frame 12035 '
        "hitting the solver's yaw and roll limits, bit it was a small matter of "
        'releasing the limits in the config, and the hard frame got solved as well."',
        6.62,
        6.22,
        6.12,
        0.98,
        font_size=10.8,
    )
    add_notes(
        slide,
        "Good beat to pause on — first real proof of concept, worth letting it breathe before complicating it.",
        visual_spec=visual_spec,
    )


def build_slide_11(prs, visual_11=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PRIMARY)
    add_title(slide, "The Other Six", dark=True, font_size=40)

    body = (
        "Two-stage auto-matching: manual correspondences as a trusted anchor, "
        "LightGlue auto-matches extending coverage outward from there. First "
        "appearance of a pattern that gets reinvented twice more before this "
        "project is done: trust a good core, extend outward carefully."
    )
    add_rect(
        slide, 1.75, 1.48, 4.45, 4.82, DEEP, "334155", radius=True, line_width=1.2
    )
    add_text(
        slide,
        body,
        2.07,
        1.84,
        3.81,
        4.10,
        18.5,
        OFFWHITE,
        font=BODY_FONT,
        valign=MSO_ANCHOR.MIDDLE,
    )

    visual_spec = (
        'A simple diagram — 7 solid "trusted" nodes in the center, 6 dashed '
        '"extended" nodes reaching out from them. [NTPvChat]'
    )
    if visual_11 is not None:
        picture = add_picture_contain(
            slide, visual_11, 6.48, 1.48, 6.32, 4.82, dark=True
        )
        picture.name = "Slide11_Visual"
    else:
        add_visual_placeholder(
            slide, visual_spec, 6.48, 1.48, 6.32, 4.82, dark=True
        )
    add_act_ii_rail(slide, 11)
    add_speech_bubble(
        slide,
        '"Yeah, so much effort, and great thinking, but it just didn\'t F***ing solve."',
        7.74,
        6.23,
        5.00,
        0.82,
        font_size=13.2,
    )
    add_notes(
        slide,
        "Plant this pattern explicitly here — it pays off big later (CLIP weighting, then the two-stage MASt3R split) and the audience should recognize it when it returns.",
        visual_spec=visual_spec,
    )


def build_slide_12(prs, visual_12_01=None, visual_12_02=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PRIMARY)
    add_title(slide, "The Mid-Way Report", dark=True, font_size=40)

    body = (
        'Built collaboratively over two late nights from ~19 screenshots: '
        "executive summary → pipeline walkthrough → what's working → a Blender "
        'comparison of solved vs. telemetry cameras ("how much the GPS is lying") '
        "→ an honest weaknesses section → a forward plan."
    )
    add_rect(
        slide, 1.75, 1.48, 4.42, 4.98, DEEP, "334155", radius=True, line_width=1.2
    )
    add_text(
        slide,
        body,
        2.08,
        1.83,
        3.76,
        4.25,
        18.4,
        OFFWHITE,
        font=BODY_FONT,
        valign=MSO_ANCHOR.MIDDLE,
    )

    visual_spec = (
        "A collage of 4–5 real screenshots from `report.html`/`img/` if they still "
        "exist on disk, or the Blender camera-comparison shot specifically — it's "
        "the most visually striking one. [NTPvCode]"
    )
    if visual_12_01 is not None:
        main_visual = add_rounded_picture_cover(
            slide, visual_12_01, 6.45, 1.48, 6.35, 4.98, line_color="516177"
        )
        main_visual.name = "Slide12_Report_Visual"
        if visual_12_02 is not None:
            add_rect(
                slide, 9.54, 4.25, 3.02, 1.94, DEEP, WHITE,
                radius=True, line_width=1.6,
            )
            inset = add_picture_cover(slide, visual_12_02, 9.62, 4.33, 2.86, 1.78)
            inset._element.spPr.prstGeom.set("prst", "roundRect")
            inset.name = "Slide12_Blender_Inset"
    else:
        add_visual_placeholder(
            slide, visual_spec, 6.45, 1.48, 6.35, 4.98, dark=True
        )
    add_inline_act_rail(slide, 12)
    add_notes(
        slide,
        "This is a good \"behind the curtain\" slide — a real deliverable, not just code, and it forces honesty about what wasn't working yet.",
        visual_spec=visual_spec,
    )


def layout_slide_13(slide, app_cam_paths=None, rigged_cam_paths=None):
    set_background(slide, OFFWHITE)
    add_title(slide, "The Honest Weakness", font_size=40)

    body = "Validating in Blender had shown that there is still a pitch problem with the solved cameras."
    add_rect(slide, 1.75, 1.25, 11.05, 0.68, WHITE, PALE, radius=True, line_width=1.2)
    add_text(
        slide, body, 2.05, 1.38, 10.45, 0.42, 16.2, PRIMARY,
        font=BODY_FONT, valign=MSO_ANCHOR.MIDDLE,
    )

    add_text(
        slide, "SOLVED CAMERAS — PITCH ERROR", 1.75, 2.05, 11.05, 0.24,
        12.2, ACCENT, font=HEADER_FONT, bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide, "AFTER SLIGHT PITCH ADJUSTMENT — ALIGNED",
        1.75, 4.30, 11.05, 0.24, 12.2, SECONDARY,
        font=HEADER_FONT, bold=True, valign=MSO_ANCHOR.MIDDLE,
    )

    frame_ids = ("04709", "04752", "10671")
    xs = (1.75, 5.51, 9.27)
    rows = (
        (app_cam_paths or (None, None, None), 2.36, ACCENT, "App"),
        (rigged_cam_paths or (None, None, None), 4.61, SECONDARY, "Rigged"),
    )
    for paths, y, color, prefix in rows:
        for frame_id, x, image_path in zip(frame_ids, xs, paths):
            if image_path is not None:
                picture = add_rounded_picture_cover(
                    slide, image_path, x, y, 3.53, 1.75,
                    line_color=color, line_width=1.4,
                )
                picture.name = f"Slide13_{prefix}_{frame_id}"
            else:
                add_visual_placeholder(
                    slide, f"{prefix.lower()}_cam_{frame_id}_in_blender.png",
                    x, y, 3.53, 1.75, dark=False,
                )
            add_text(
                slide, frame_id.lstrip("0"), x + 0.10, y + 0.08,
                0.68, 0.24, 10.0, WHITE, font=HEADER_FONT,
                bold=True, valign=MSO_ANCHOR.MIDDLE,
            )


def build_slide_13(prs, app_cam_paths=None, rigged_cam_paths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    layout_slide_13(slide, app_cam_paths, rigged_cam_paths)
    add_inline_act_rail(slide, 13)
    add_notes(
        slide,
        "Even though the raycasting in the app seemed good enough, validation in Blender showed that something still wasn't solving correctly.",
    )


def build_slide_14(prs, visual_14=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PRIMARY)
    add_title(slide, '"The Rabbit Hole in the Rabbit Hole"', dark=True, font_size=36)

    body = (
        "The mid-way report wasn't just a status update — it was built to ask a "
        "real question of the reviewing team: hit a wall with the 6 hard frames, "
        "considering two very different moves — upgrade the feature matcher "
        "(LoFTR/RoMa), or go build a synthetic-data YOLO skeleton network to "
        "precisely locate van joints. Which one, or neither?"
    )
    add_rect(
        slide, 1.75, 1.48, 5.02, 4.98, DEEP, "334155", radius=True, line_width=1.2
    )
    add_text(
        slide,
        body,
        2.07,
        1.78,
        4.38,
        4.32,
        17.2,
        OFFWHITE,
        font=BODY_FONT,
        valign=MSO_ANCHOR.MIDDLE,
    )

    visual_spec = '"The Rabbit Hole in the Rabbit Hole" [NTPvChat]'
    if visual_14 is not None:
        picture = add_rounded_picture_cover(
            slide, visual_14, 7.04, 1.48, 5.76, 4.98, line_color="516177"
        )
        picture.name = "Slide14_RabbitHole_Visual"
    else:
        add_visual_placeholder(
            slide, visual_spec, 7.04, 1.48, 5.76, 4.98, dark=True
        )
    add_inline_act_rail(slide, 14)
    add_speech_bubble(
        slide,
        '"I know the guy - without guidance he\'d go train a YOLO8 network. Good thing he knows himself as well."',
        6.94,
        6.12,
        5.80,
        0.94,
        font_size=11.8,
    )
    add_notes(
        slide,
        "No punchline here — play this one straight. The next slide is the answer, and it lands harder if this one isn't already winking at the audience.",
        visual_spec=visual_spec,
    )


def refresh_act_iii(prs, visual_12_01, visual_12_02, app_cam_paths, rigged_cam_paths, visual_14):
    """Refresh Slides 12–14 without touching any earlier slide."""
    if len(prs.slides) != 14:
        raise ValueError(f"Expected 14 slides, found {len(prs.slides)}.")

    slide_12 = prs.slides[11]
    for shape in list(slide_12.shapes):
        if shape.left >= Inches(6.35) and Inches(1.30) <= shape.top < Inches(6.60):
            shape._element.getparent().remove(shape._element)
    main_visual = add_rounded_picture_cover(
        slide_12, visual_12_01, 6.45, 1.48, 6.35, 4.98, line_color="516177"
    )
    main_visual.name = "Slide12_Report_Visual"
    add_rect(
        slide_12, 9.54, 4.25, 3.02, 1.94, DEEP, WHITE,
        radius=True, line_width=1.6,
    )
    inset = add_picture_cover(slide_12, visual_12_02, 9.62, 4.33, 2.86, 1.78)
    inset._element.spPr.prstGeom.set("prst", "roundRect")
    inset.name = "Slide12_Blender_Inset"

    slide_13 = prs.slides[12]
    for shape in list(slide_13.shapes):
        if shape.left >= Inches(1.43):
            shape._element.getparent().remove(shape._element)
    layout_slide_13(slide_13, app_cam_paths, rigged_cam_paths)
    add_notes(
        slide_13,
        "Even though the raycasting in the app seemed good enough, validation in Blender showed that something still wasn't solving correctly.",
    )

    slide_14 = prs.slides[13]
    for shape in list(slide_14.shapes):
        if (
            shape.name != "SpeechOverlay_Group"
            and shape.left >= Inches(6.80)
            and Inches(1.30) <= shape.top < Inches(6.60)
        ):
            shape._element.getparent().remove(shape._element)
    rabbit = add_rounded_picture_cover(
        slide_14, visual_14, 7.04, 1.48, 5.76, 4.98, line_color="516177"
    )
    rabbit.name = "Slide14_RabbitHole_Visual"
    bring_speech_overlays_to_front(prs)


def build_slide_15(prs, visual_15=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PRIMARY)
    add_title(slide, "The Answers", dark=True, font_size=40)

    cards = (
        (
            1.48,
            2.20,
            "Vector 1 — Let us save you from yourself: DROP the custom YOLO idea. Instead, pair an open-dictionary vision-language model (Qwen) with CLIP embeddings to dynamically discover and correlate geometric anchors, no predefined target needed.",
            SECONDARY,
        ),
        (
            3.86,
            1.62,
            "Vector 2 — Yes, upgrade the solve engine, but use MASt3R, which outputs dense 3D maps and local features jointly.",
            SECONDARY,
        ),
    )
    for y, h, text, color in cards:
        add_rect(slide, 1.75, y, 5.92, h, DEEP, "334155", radius=True, line_width=1.1)
        add_text(
            slide, text, 2.05, y + 0.20, 5.32, h - 0.40, 15.2,
            OFFWHITE, font=BODY_FONT, valign=MSO_ANCHOR.MIDDLE,
        )
        marker = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.75), Inches(y), Inches(0.065), Inches(h),
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = rgb(color)
        marker.line.fill.background()

    add_rect(slide, 1.75, 5.66, 5.92, 0.80, "182338", SECONDARY, radius=True, line_width=1.0)
    add_text(
        slide,
        "Both technologies that show up in the codebase exactly six days later were named, specifically, in this document.",
        2.03, 5.81, 5.36, 0.48, 13.2, SECONDARY,
        font=HEADER_FONT, bold=True, valign=MSO_ANCHOR.MIDDLE,
    )

    visual_spec = (
        'Cartoon image: YOLO8 character laying on the ground unconscious, A drone titled "Xtend" with a spring boxing glove unhinged (it had just knocked out YOLO8 character), me looking at the drone in a "oh no, don\'t hurt YOLO8 character!" pose. [NTPvChat]'
    )
    if visual_15 is not None:
        picture = add_picture_contain(slide, visual_15, 7.95, 1.48, 4.85, 4.98, dark=True)
        picture.name = "Slide15_Visual"
    else:
        add_visual_placeholder(slide, visual_spec, 7.95, 1.48, 4.85, 4.98, dark=True)
    add_inline_act_rail(slide, 15)
    add_notes(
        slide,
        "Worth naming plainly: this wasn't an independent rethink arrived at in isolation. It was a direct, specific technical directive, followed closely. That's a more honest and more interesting story than \"we had a eureka moment\" — say so. Also plant this, quietly, for later: the feedback doc suggests a CLIP cosine-similarity cutoff of >0.85. Remember that number.",
        visual_spec=visual_spec,
    )


def build_slide_16(prs, visual_16=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, OFFWHITE)
    add_title(slide, "The Pivot We Said No To", font_size=40)

    body = (
        "A follow-up call raised a third option: move to continuous video, reconstructed by MASt3R, potentially easing multi-drone alignment altogether. Declined - on basis of: The hardware on hand couldn't support it, lack of data, and some research alerting MASt3R drifting, without a good anchor, Qwen could supply. Rather than chase a second rabbit hole while still climbing out of the first, the project stayed the course."
    )
    add_rect(slide, 1.75, 1.48, 5.02, 4.98, WHITE, PALE, radius=True, line_width=1.2)
    add_text(
        slide, body, 2.07, 1.78, 4.38, 4.32, 16.5, PRIMARY,
        font=BODY_FONT, valign=MSO_ANCHOR.MIDDLE,
    )
    visual_spec = (
        'A simple fork-in-the-road diagram: one path taken (labeled "frames + MASt3R + Qwen"), one path greyed out and clearly signposted "not this time" (labeled "continuous video reconstruction"). [NTPvChat]'
    )
    if visual_16 is not None:
        picture = add_picture_contain(slide, visual_16, 7.04, 1.48, 5.76, 4.98, dark=False)
        picture.name = "Slide16_Visual"
    else:
        add_visual_placeholder(slide, visual_spec, 7.04, 1.48, 5.76, 4.98, dark=False)
    add_inline_act_rail(slide, 16)
    add_notes(
        slide,
        "This is a scope-discipline beat, not a modesty beat — worth saying plainly that turning down a good idea, on purpose, because of real constraints, is its own kind of engineering judgment.",
        visual_spec=visual_spec,
    )


def build_slide_17(prs, video_17=None, poster_17=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PRIMARY)
    add_text(
        slide, "Deep", 1.75, 0.36, 1.42, 0.62, 40, ACCENT,
        font=HEADER_FONT, bold=True, valign=MSO_ANCHOR.MIDDLE,
    )
    strike = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.74), Inches(0.68), Inches(1.06), Inches(0.045)
    )
    strike.fill.solid()
    strike.fill.fore_color.rgb = rgb(ACCENT)
    strike.line.fill.background()
    add_text(
        slide, "Fast Learning", 3.18, 0.36, 4.00, 0.62, 40, WHITE,
        font=HEADER_FONT, bold=True, valign=MSO_ANCHOR.MIDDLE,
    )
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.75), Inches(1.10), Inches(0.78), Inches(0.055)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(SECONDARY)
    accent.line.fill.background()

    body = (
        "Report sent May 31st. Feedback landed June 8th (previous two slides). Between the answer arriving and development actually resuming June 14th: a week of DeepLearning.AI coursework — Agentic AI, and a Claude-specific course."
    )
    add_rect(slide, 1.75, 1.48, 5.02, 4.98, DEEP, "334155", radius=True, line_width=1.2)
    add_text(
        slide, body, 2.07, 1.78, 4.38, 4.32, 18.0, OFFWHITE,
        font=BODY_FONT, valign=MSO_ANCHOR.MIDDLE,
    )
    visual_spec = 'Maybe the "I know Code Too" meme. [NTPvCode]'
    if video_17 is not None and poster_17 is not None:
        add_movie_contain(
            slide, video_17, poster_17, 7.04, 1.48, 5.76, 4.98,
            "Slide17_Video", dark=True,
        )
    else:
        add_visual_placeholder(slide, visual_spec, 7.04, 1.48, 5.76, 4.98, dark=True)
    add_inline_act_rail(slide, 17)
    add_notes(
        slide,
        "Frame this as \"here's what filled the rest of the visible gap,\" not as the reveal itself — the reveal already happened two slides ago.",
        visual_spec=visual_spec,
    )


def build_slide_18(prs, visual_18=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, OFFWHITE)
    add_title(slide, "Tools, Tokens and Model Diversity", font_size=36)

    body = (
        "Up until now, work got done in Claude's web chat. Files were moved manually, which felt frugal but measurably wasn't. Same day as the pivot commit, development moved to Claude Code: more output, less friction, immediately. Second habit, formed alongside it: plan in Claude Code, get a second opinion from Gemini, feed that back into implementation — cross-AI review that caught logic bugs early and kept plan-level questions off Claude's context."
    )
    add_rect(slide, 1.75, 1.48, 5.02, 4.98, WHITE, PALE, radius=True, line_width=1.2)
    add_text(
        slide, body, 2.07, 1.75, 4.38, 4.44, 15.8, PRIMARY,
        font=BODY_FONT, valign=MSO_ANCHOR.MIDDLE,
    )
    visual_spec = (
        'Two small workflow loops side by side. Before: Human ↔ (manual upload/download) ↔ Claude web chat, drawn with visible friction (a little "⏳" or paperclip icon on the transfer arrows). After: a triangle — Human/Claude Code plan → Gemini review → feedback back into Claude Code for implementation. [NTPvChat]'
    )
    if visual_18 is not None:
        picture = add_picture_contain(slide, visual_18, 7.04, 1.48, 5.76, 4.98, dark=False)
        picture.name = "Slide18_Visual"
    else:
        add_visual_placeholder(slide, visual_spec, 7.04, 1.48, 5.76, 4.98, dark=False)
    add_inline_act_rail(slide, 18)
    add_notes(
        slide,
        "Not a confession — a measured course-correction, same honesty the rest of this deck runs on: an assumed efficiency that turned out to be a false economy once actually checked against outcomes, not intuition. Good beat to close the act on, since it's the actual last thing that changed before the pivot itself.",
        visual_spec=visual_spec,
    )


def build_slide_19(prs, visual_19=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PRIMARY)
    add_title(slide, "Heart Surgery", dark=True, font_size=40)

    body = (
        "Deleted outright: LightGlue, GroundingDINO, SAM, SuperPoint, the debug UI built for them. In: MASt3R-SfM (dense 3D reconstruction + camera poses) + Qwen VL (open-world object discovery) + CLIP (confidence weighting)."
    )
    add_rect(slide, 1.75, 1.48, 5.02, 4.98, DEEP, "334155", radius=True, line_width=1.2)
    add_text(
        slide, body, 2.07, 1.78, 4.38, 4.32, 17.0, OFFWHITE,
        font=BODY_FONT, valign=MSO_ANCHOR.MIDDLE,
    )
    visual_spec = (
        'A stark "DELETED" vs "ADDED" two-column list, red strikethrough vs. green. [NTPvChat]'
    )
    if visual_19 is not None:
        picture = add_picture_contain(slide, visual_19, 7.04, 1.48, 5.76, 4.98, dark=True)
        picture.name = "Slide19_Visual"
    else:
        add_visual_placeholder(slide, visual_spec, 7.04, 1.48, 5.76, 4.98, dark=True)
    add_inline_act_rail(slide, 19, ACT_IV_STATIONS)
    add_notes(
        slide,
        "Don't soften this one — the scale of the rewrite is the point. Everything code-related from Act II except the OCR/GPS core and `camera_deltas.py` is gone. (Act III was never code to begin with — it's the human decision that triggered this slide.)",
        visual_spec=visual_spec,
    )


def layout_slide_20(slide, visual_20=None):
    set_background(slide, OFFWHITE)
    add_title(slide, "Enter MASt3R", font_size=40)

    body = (
        "The old stack just couldn't crack it, and something smarter was needed. MASt3R's dense reconstruction could handle the low-texture, repetitive ground the old feature matchers choked on."
    )
    add_rect(slide, 1.75, 1.48, 5.02, 4.98, WHITE, PALE, radius=True, line_width=1.2)
    add_text(
        slide, body, 2.07, 1.78, 4.38, 4.32, 19.0, PRIMARY,
        font=BODY_FONT, valign=MSO_ANCHOR.MIDDLE,
    )
    visual_spec = "Created."
    if visual_20 is not None:
        picture = add_picture_contain(slide, visual_20, 7.04, 1.48, 5.76, 4.98, dark=False)
        picture.name = "Slide20_Visual"
    else:
        add_visual_placeholder(slide, visual_spec, 7.04, 1.48, 5.76, 4.98, dark=False)
    add_notes(
        slide,
        "Mention that this magic requires better hardware.",
        visual_spec=visual_spec,
    )


def build_slide_20(prs, visual_20=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    layout_slide_20(slide, visual_20)
    add_inline_act_rail(slide, 20, ACT_IV_STATIONS)


def refresh_slide_20(prs, visual_20=None):
    """Replace only Slide 20's approved content, preserving its Act IV rail."""
    if len(prs.slides) != 22:
        raise ValueError(f"Expected 22 slides, found {len(prs.slides)}.")
    slide = prs.slides[19]
    for shape in list(slide.shapes):
        if shape.left >= Inches(1.43):
            shape._element.getparent().remove(shape._element)
    layout_slide_20(slide, visual_20)


def build_slide_21(prs, visual_21=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PRIMARY)
    add_title(slide, "Step Zero: Get a GPU", dark=True, font_size=40)

    body = (
        "RunPod, chosen over Google Colab specifically because Colab's runtimes wipe on timeout — meaning MASt3R's CUDA/C++ kernels would need recompiling from scratch every session, versus RunPod's persistent volume storage compiling once and staying ready."
    )
    add_rect(slide, 1.75, 1.48, 5.02, 4.98, DEEP, "334155", radius=True, line_width=1.2)
    add_text(
        slide, body, 2.07, 1.78, 4.38, 4.32, 17.0, OFFWHITE,
        font=BODY_FONT, valign=MSO_ANCHOR.MIDDLE,
    )
    visual_spec = (
        'A simple two-box comparison card, Colab vs. RunPod, with the "recompile every session" vs "compile once" distinction as the headline. [NTPvChat]'
    )
    if visual_21 is not None:
        picture = add_picture_contain(slide, visual_21, 7.04, 1.48, 5.76, 4.98, dark=True)
        picture.name = "Slide21_Visual"
    else:
        add_visual_placeholder(slide, visual_spec, 7.04, 1.48, 5.76, 4.98, dark=True)
    add_inline_act_rail(slide, 21, ACT_IV_STATIONS)
    add_notes(
        slide,
        "Frame this as an actual evaluated decision, not a default — it's easy to skip over infrastructure choices in a deck, don't.",
        visual_spec=visual_spec,
    )


def build_slide_23(prs, visual_23=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, OFFWHITE)
    add_title(slide, "Compute Headless in the Cloud, Test Locally", font_size=35)

    body_segments = (
        ("RunPod has no display; the interactive viewer needs one. ", BODY_FONT),
        ("--export-solve", "Cascadia Mono"),
        (" runs the full pipeline remotely and serializes camera poses + terrain point clouds to disk. Those files travel to the local machine, where ", BODY_FONT),
        ("--import-solve", "Cascadia Mono"),
        (" skips OCR/MASt3R/Ceres entirely and opens the viewer straight against the already-solved cameras.", BODY_FONT),
    )
    add_rect(slide, 1.75, 1.48, 5.02, 4.98, WHITE, PALE, radius=True, line_width=1.2)
    add_rich_text(
        slide, body_segments, 2.07, 1.72, 4.38, 4.48, 15.8, PRIMARY,
        valign=MSO_ANCHOR.MIDDLE,
    )
    visual_spec = (
        'A literal "cloud → local" diagram: RunPod box producing `solved_cameras.json` + point-cloud files, an arrow labeled "copy down," a local laptop box opening the viewer. [NTPvChat]'
    )
    if visual_23 is not None:
        picture = add_picture_contain(slide, visual_23, 7.04, 1.48, 5.76, 4.98, dark=False)
        picture.name = "Slide23_Visual"
    else:
        add_visual_placeholder(slide, visual_spec, 7.04, 1.48, 5.76, 4.98, dark=False)
    add_inline_act_rail(slide, 23, ACT_IV_STATIONS)
    add_notes(
        slide,
        "Close this slide on the line that sets up the next act: for most of this project, the only thing crossing that gap was a handful of numbers — the actual 3D reconstruction never left the pod. That's about to matter.",
        visual_spec=visual_spec,
    )


def build_slide_23_5(prs, visual_23_5=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PRIMARY)
    add_title(slide, "MASt3R's Achilles' heel", dark=True, font_size=38)

    body = (
        "MASt3R has one weakness - it DRIFTS. We will use QWEN to find a joint anchor in the images, to prevent this."
    )
    add_rect(slide, 1.75, 1.48, 5.02, 4.98, DEEP, "334155", radius=True, line_width=1.2)
    add_text(
        slide, body, 2.07, 1.78, 4.38, 4.32, 20.0, OFFWHITE,
        font=BODY_FONT, valign=MSO_ANCHOR.MIDDLE,
    )
    visual_spec = "Created."
    if visual_23_5 is not None:
        picture = add_picture_contain(
            slide, visual_23_5, 7.04, 1.48, 5.76, 4.98, dark=True
        )
        picture.name = "Slide23_5_Visual"
    else:
        add_visual_placeholder(slide, visual_spec, 7.04, 1.48, 5.76, 4.98, dark=True)
    add_inline_act_rail(slide, 23.5, ACT_V_STATIONS)
    add_notes(
        slide,
        "Next slide shows the method in more detail.",
        visual_spec=visual_spec,
    )


def build_slide_24(prs, visual_24=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, OFFWHITE)
    add_title(slide, "Qwen's Job: Find the One Thing Every Frame Shares", font_size=33)

    body = (
        "Per frame: list every man-made object, bbox + label. Across all frames: consolidate labels, rank by how many frames contain each one. Among the top-coverage candidates, a second Qwen call — reasoning over label text alone — picks the single best fixed 3D reference point. Output: that object's pixel centroid in every frame it appears in, weighted by CLIP confidence."
    )
    add_rect(slide, 1.75, 1.48, 5.02, 4.98, WHITE, PALE, radius=True, line_width=1.2)
    add_text(
        slide, body, 2.07, 1.70, 4.38, 4.52, 15.4, PRIMARY,
        font=BODY_FONT, valign=MSO_ANCHOR.MIDDLE,
    )
    visual_spec = (
        'A 3-step flow: "find everything" → "which label shows up everywhere" → "which of those is trustworthy" — with the van as the visual through-line landing at the end. [NTPvChat]'
    )
    if visual_24 is not None:
        picture = add_picture_contain(slide, visual_24, 7.04, 1.48, 5.76, 4.98, dark=False)
        picture.name = "Slide24_Visual"
    else:
        add_visual_placeholder(slide, visual_spec, 7.04, 1.48, 5.76, 4.98, dark=False)
    add_inline_act_rail(slide, 24, ACT_V_STATIONS)
    add_notes(
        slide,
        'Emphasize this is open-world — nothing here hardcodes "find the van." It discovers that the van is the answer.',
        visual_spec=visual_spec,
    )


def build_slide_25(prs, anchor_images, mascot_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PRIMARY)
    add_title(slide, "Debug-By-Drawing", dark=True, font_size=40)

    body_segments = (
        ("Initially, the code was reading QWEN's data wrong, and the Anchor Centroids were offset badly. So we added the ", BODY_FONT),
        ("--preview-anchor", "Cascadia Mono"),
        (" arg to get some debug drawing in order to figure out what's going on.", BODY_FONT),
    )
    add_rect(slide, 1.75, 1.48, 4.15, 4.98, DEEP, "334155", radius=True, line_width=1.2)
    add_rich_text(
        slide, body_segments, 2.07, 1.76, 3.51, 4.40, 17.0, OFFWHITE,
        valign=MSO_ANCHOR.MIDDLE,
    )

    visual_spec = 'a few of the images in the "anchor" folder'
    add_rect(slide, 6.15, 1.48, 6.65, 4.98, "182338", "516177", radius=True, line_width=1.2)
    slots = (
        (6.30, 1.83),
        (9.60, 1.83),
        (6.30, 3.84),
        (9.60, 3.84),
    )
    for index, ((x, y), image_path) in enumerate(zip(slots, anchor_images), start=1):
        picture = add_rounded_picture_cover(slide, image_path, x, y, 3.05, 1.72, line_color="516177")
        picture.name = f"Slide25_Anchor_{index}"

    add_inline_act_rail(slide, 25, ACT_V_STATIONS)
    bubble_text = (
        "...and to presuade Claude that there really IS a coordinate system bug, because it insisted it's a hundred other things."
    )
    add_speech_bubble(slide, bubble_text, 6.32, 4.72, 5.28, 1.05, font_size=13.0)
    mascot = slide.shapes.add_picture(
        str(mascot_path),
        Inches(MASCOT_X),
        Inches(MASCOT_Y),
        Inches(MASCOT_SIZE),
        Inches(MASCOT_SIZE),
    )
    mascot.name = "Mascot"
    configure_speech_overlay_click_timing(prs, 25)
    add_notes(slide, "", visual_spec=visual_spec)


def refresh_slide_25_callout(prs, mascot_path):
    """Replace Slide 25's legacy composite bubble with one native callout."""
    if len(prs.slides) != 25:
        raise ValueError(f"Expected 25 slides, found {len(prs.slides)}.")
    slide = prs.slides[24]
    legacy_names = {
        "SpeechOverlay_Group",
        "Mascot_Pivot_Frame",
        "SpeechBubble_Background",
        "SpeechBubble_Text",
        "SpeechBubble_Tail",
        "SpeechBubble_Callout",
        "Mascot",
    }
    for shape in list(slide.shapes):
        if shape.name in legacy_names:
            shape._element.getparent().remove(shape._element)

    bubble_text = (
        "...and to presuade Claude that there really IS a coordinate system bug, because it insisted it's a hundred other things."
    )
    add_speech_bubble(slide, bubble_text, 6.32, 4.72, 5.28, 1.05, font_size=13.0)
    mascot = slide.shapes.add_picture(
        str(mascot_path),
        Inches(MASCOT_X),
        Inches(MASCOT_Y),
        Inches(MASCOT_SIZE),
        Inches(MASCOT_SIZE),
    )
    mascot.name = "Mascot"
    configure_speech_overlay_click_timing(prs, 25)


def build_slide_26(prs, visual_26=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PRIMARY)
    add_title(slide, "New Engines - New Bugs", dark=True, font_size=40)

    items = (
        "Wrong CDN paths (MASt3R's weights aren't available on Hugging Face, so had to be downloaded from Naver's CDN)",
        "Trying to read 'dust3r' attributes in 'MASt3R' class",
        "'CLIP_ANCHOR_THRESHOLD' too high (QWEN)",
        "Wrong coordinate systems for QWEN (mentioned in previous slide), and for MASt3R's camera solve.",
    )
    slots = (
        (1.75, 1.48),
        (4.35, 1.48),
        (1.75, 3.98),
        (4.35, 3.98),
    )
    for index, (text, (x, y)) in enumerate(zip(items, slots), start=1):
        item_font_size = 10.8 if index == 3 else 12.2
        add_rect(slide, x, y, 2.42, 2.30, DEEP, ACCENT, radius=True, line_width=1.1)
        add_text(
            slide, f"{index:02d}", x + 0.17, y + 0.14, 0.38, 0.24,
            10, ACCENT, font=HEADER_FONT, bold=True,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide, text, x + 0.17, y + 0.48, 2.08, 1.61,
            item_font_size, OFFWHITE, font=BODY_FONT,
            valign=MSO_ANCHOR.MIDDLE,
        )

    if visual_26 is not None:
        picture = add_picture_contain(slide, visual_26, 7.04, 1.48, 5.76, 4.98, dark=True)
        picture.name = "Slide26_Visual"
    else:
        add_visual_placeholder(slide, "", 7.04, 1.48, 5.76, 4.98, dark=True)
    add_inline_act_rail(slide, 26)
    add_notes(slide, "")


def build_slide_27(prs, visual_27=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, OFFWHITE)
    add_title(slide, "QWEN isn't a Psychic", font_size=40)

    body = (
        "In order for MASt3R to give best results we had to make sure the frames are clean of overlay graphics. The last piece of overlay graphics we needed to take care of, was the moving crosshair. We wanted to use QWEN for that, but had to describe the crosshair better:\n"
        'Rather than: "a crosshair symbol", "+" or "⊕ symbol"\n'
        'When used: "a small ring of four chevron marks around a center dot" QWEN found it and we blurred it out.'
    )
    add_rect(slide, 1.75, 1.48, 5.02, 4.98, WHITE, PALE, radius=True, line_width=1.2)
    add_text(
        slide, body, 2.07, 1.70, 4.38, 4.52, 15.0, PRIMARY,
        font=BODY_FONT, valign=MSO_ANCHOR.MIDDLE,
    )
    visual_spec = (
        'Side-by-side: what was described (a big fat "+") vs. what\'s actually on screen (the chevron-ring reticle), with the wrong guessed location circled. [NTPvChat]'
    )
    if visual_27 is not None:
        picture = add_picture_contain(slide, visual_27, 7.04, 1.48, 5.76, 4.98, dark=False)
        picture.name = "Slide27_Visual"
    else:
        add_visual_placeholder(slide, visual_spec, 7.04, 1.48, 5.76, 4.98, dark=False)
    add_inline_act_rail(slide, 27)
    add_notes(slide, "", visual_spec=visual_spec)


def build_slide_28(prs, video_28=None, poster_28=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PRIMARY)
    add_text(
        slide, "--export-mesh", 1.75, 0.36, 4.05, 0.62, 36, WHITE,
        font="Cascadia Mono", bold=True, valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide, ": Actually Looking", 5.54, 0.36, 5.94, 0.62, 36, WHITE,
        font=HEADER_FONT, bold=True, valign=MSO_ANCHOR.MIDDLE,
    )
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.75), Inches(1.10), Inches(0.78), Inches(0.055)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(SECONDARY)
    accent.line.fill.background()

    body = (
        "Even after cleaning the input frames perfectly, the reconstruction wasn't good. So it was time to see things with our eyes: a debug tool that dumps the raw MASt3R reconstruction + solved cameras (and, once solved, the Ceres-refined + Sim(3)-aligned version) as real .glb scenes, importable straight into Blender."
    )
    add_rect(slide, 1.75, 1.48, 5.02, 4.98, DEEP, "334155", radius=True, line_width=1.2)
    add_text(
        slide, body, 2.07, 1.72, 4.38, 4.50, 16.0, OFFWHITE,
        font=BODY_FONT, valign=MSO_ANCHOR.MIDDLE,
    )
    visual_spec = (
        "A real Blender screenshot of the point cloud + camera cards, if available; otherwise a simple render mockup of scattered 3D points with camera frustums. [NTPvCode]"
    )
    if video_28 is not None and poster_28 is not None:
        add_movie_contain(
            slide, video_28, poster_28, 7.04, 1.48, 5.76, 4.98,
            "Slide28_Video", dark=True,
        )
    else:
        add_visual_placeholder(slide, visual_spec, 7.04, 1.48, 5.76, 4.98, dark=True)
    add_inline_act_rail(slide, 28)
    add_notes(slide, "", visual_spec=visual_spec)


def refresh_slides_15_18(prs, visual_15, visual_16, video_17, poster_17, visual_18):
    """Replace Slides 15–18 visual placeholders without changing their text."""
    if len(prs.slides) != 18:
        raise ValueError(f"Expected 18 slides, found {len(prs.slides)}.")

    specs = (
        (15, 7.85, visual_15, "picture", True, "Slide15_Visual"),
        (16, 6.94, visual_16, "picture", False, "Slide16_Visual"),
        (17, 6.94, video_17, "movie", True, "Slide17_Video"),
        (18, 6.94, visual_18, "picture", False, "Slide18_Visual"),
    )
    for slide_number, threshold, asset, kind, dark, shape_name in specs:
        slide = prs.slides[slide_number - 1]
        for shape in list(slide.shapes):
            if shape.left >= Inches(threshold) and Inches(1.30) <= shape.top < Inches(6.60):
                shape._element.getparent().remove(shape._element)

        x = 7.95 if slide_number == 15 else 7.04
        w = 4.85 if slide_number == 15 else 5.76
        if kind == "movie":
            add_movie_contain(
                slide, asset, poster_17, x, 1.48, w, 4.98,
                shape_name, dark=dark,
            )
        else:
            picture = add_picture_contain(slide, asset, x, 1.48, w, 4.98, dark=dark)
            picture.name = shape_name


def normalize_slide_10_title(prs):
    """Keep the long Slide 10 title clear of the accent rule in existing decks."""
    if len(prs.slides) < 10:
        return
    expected_title = '"Manual Feature Matching Works Great on 7 Frames"'
    for shape in prs.slides[9].shapes:
        if getattr(shape, "has_text_frame", False) and shape.text == expected_title:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(28)
            return


def main():
    parser = argparse.ArgumentParser(
        description="Build the Raycast Challenge deck through Act VI, skipping Slide 22."
    )
    parser.add_argument(
        "--visuals-dir",
        type=Path,
        default=DEFAULT_VISUALS_DIR,
        help="Folder containing generated slide visuals.",
    )
    parser.add_argument(
        "--dataset-dir",
        type=Path,
        default=DEFAULT_DATASET_DIR,
        help="Folder containing source dataset images.",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=OUTPUT_PPTX,
        help="Destination .pptx path.",
    )
    parser.add_argument(
        "--mascot",
        type=Path,
        default=DEFAULT_MASCOT_PATH,
        help="Transparent mascot PNG used by speech overlays.",
    )
    parser.add_argument(
        "--serious-mascot",
        type=Path,
        default=DEFAULT_SERIOUS_MASCOT_PATH,
        help="Red-portal mascot PNG used on Slide 11 only.",
    )
    parser.add_argument(
        "--skip-animations",
        action="store_true",
        help="Save without applying PowerPoint entrance animations.",
    )
    parser.add_argument(
        "--refresh-act-iii",
        action="store_true",
        help="Refresh only Slides 12–14 in an existing 14-slide deck.",
    )
    parser.add_argument(
        "--refresh-slides-15-18",
        action="store_true",
        help="Replace only the visual placeholders on Slides 15–18.",
    )
    parser.add_argument(
        "--refresh-slide-20",
        action="store_true",
        help="Replace only Slide 20 with its approved rewrite.",
    )
    parser.add_argument(
        "--refresh-slide-25-callout",
        action="store_true",
        help="Replace Slide 25's speech overlay with a native PowerPoint callout.",
    )
    args = parser.parse_args()

    visual_1 = find_asset(args.visuals_dir, "Raycast_Slide_1_Visual.png")
    visual_2 = find_asset(args.visuals_dir, "Raycast_Slide_2_Visual.png")
    visual_4 = find_asset(args.visuals_dir, "Raycast_Slide_4_Visual.png")
    visual_5 = find_asset(args.visuals_dir, "Raycast_Slide_5_Visual.png")
    visual_6 = find_asset(args.visuals_dir, "Raycast_Slide_6_Visual.png")
    visual_7 = find_asset(args.visuals_dir, "Raycast_Slide_7_Visual.png")
    visual_8 = find_asset(args.visuals_dir, "Raycast_Slide_8_Visual_H264_4K.mp4")
    poster_8 = find_asset(args.visuals_dir, "Raycast_Slide_8_Poster.jpg")
    visual_9 = find_asset(args.visuals_dir, "Raycast_Slide_9_Visual.png")
    visual_10 = find_asset(
        args.visuals_dir, "Raycast_Slide_10_Visual_H264_4K.mp4"
    )
    poster_10 = find_asset(args.visuals_dir, "Raycast_Slide_10_Poster.jpg")
    visual_11 = find_asset(args.visuals_dir, "Raycast_Slide_11_Visual.png")
    visual_12_01 = find_asset(args.visuals_dir, "Raycast_Slide_12_Visual_01.png")
    visual_12_02 = find_asset(args.visuals_dir, "Raycast_Slide_12_Visual_02.png")
    visual_14 = find_asset(args.visuals_dir, "Raycast_Slide_14_Visual.png")
    visual_15 = find_asset(args.visuals_dir, "Raycast_Slide_15_Visual.png")
    visual_16 = find_asset(args.visuals_dir, "Raycast_Slide_16_Visual.png")
    video_17 = find_asset(args.visuals_dir, "Raycast_Slide_17_Visual.mp4")
    poster_17 = find_asset(args.visuals_dir, "Raycast_Slide_17_Poster.jpg")
    visual_18 = find_asset(args.visuals_dir, "Raycast_Slide_18_Visual.png")
    visual_19 = find_asset(args.visuals_dir, "Raycast_Slide_19_Visual.png")
    visual_20 = find_asset(args.visuals_dir, "Raycast_Slide_20_NEW_Visual.png")
    visual_21 = find_asset(args.visuals_dir, "Raycast_Slide_21_Visual.png")
    visual_23 = find_asset(args.visuals_dir, "Raycast_Slide_23_Visual.png")
    visual_23_5 = find_asset(args.visuals_dir, "Raycast_Slide_23.5_Visual.png")
    visual_24 = find_asset(args.visuals_dir, "Raycast_Slide_24_Visual.png")
    visual_26 = find_asset(args.visuals_dir, "Raycast_Slide_26_Visual.png")
    visual_27 = find_asset(args.visuals_dir, "Raycast_Slide_27_Visual.png")
    video_28 = find_asset(args.visuals_dir, "Raycast_Slide_28_Visual.mp4")
    poster_28 = find_asset(args.visuals_dir, "Raycast_Slide_28_Poster.jpg")
    anchor_dir = args.visuals_dir / "anchor"
    anchor_images = tuple(
        find_asset(anchor_dir, filename)
        for filename in (
            "2026-02-15_16-25-03_04569_anchor.jpg",
            "2026-02-15_16-25-03_04709_anchor.jpg",
            "2026-02-15_16-25-03_10671_anchor.jpg",
            "2026-02-15_16-35-56_09763_anchor.jpg",
        )
    )
    frame_ids = ("04709", "04752", "10671")
    app_cam_paths = tuple(
        find_asset(args.visuals_dir, f"app_cam_{frame_id}_in_blender.png")
        for frame_id in frame_ids
    )
    rigged_cam_paths = tuple(
        find_asset(args.visuals_dir, f"rigged_cam_{frame_id}_in_blender.png")
        for frame_id in frame_ids
    )
    sharp_frame = find_asset(args.dataset_dir, "2026-02-15_16-25-03_12035.png")
    blurred_frame = find_asset(args.dataset_dir, "2026-02-15_16-25-03_04681.png")

    args.output.parent.mkdir(parents=True, exist_ok=True)
    if args.output.exists():
        prs = Presentation(str(args.output))
        if len(prs.slides) not in (3, 6, 11, 14, 18, 22, 25, 28):
            raise ValueError(
                f"Expected 3, 6, 11, 14, 18, 22, 25, or 28 existing slides in {args.output}, "
                f"found {len(prs.slides)}."
            )
        starting_slide_count = len(prs.slides)
    else:
        prs = Presentation()
        starting_slide_count = 0
        prs.slide_width = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)
        prs.core_properties.title = "The Raycast Challenge — Acts 0 through II"
        prs.core_properties.subject = "The Setup, Bring-Up, and By Hand"
        prs.core_properties.author = "Ofer Shkedi"
        prs.core_properties.keywords = "Raycast, drone, reprojection, calibration"

        build_slide_1(prs, visual_1)
        build_slide_2(prs, visual_2)
        build_slide_3(prs, sharp_frame, blurred_frame)

    if args.refresh_act_iii:
        refresh_act_iii(
            prs,
            visual_12_01,
            visual_12_02,
            app_cam_paths,
            rigged_cam_paths,
            visual_14,
        )
        prs.save(args.output)
        print(args.output)
        return

    if args.refresh_slide_20:
        refresh_slide_20(prs, visual_20)
        prs.save(args.output)
        print(args.output)
        return

    if args.refresh_slide_25_callout:
        refresh_slide_25_callout(prs, args.mascot)
        prs.save(args.output)
        print(args.output)
        return

    if args.refresh_slides_15_18:
        refresh_slides_15_18(
            prs, visual_15, visual_16, video_17, poster_17, visual_18
        )
        prs.save(args.output)
        configure_slide17_video(args.output)
        print(args.output)
        return

    if starting_slide_count == 14:
        build_slide_15(prs, visual_15)
        build_slide_16(prs, visual_16)
        build_slide_17(prs, video_17, poster_17)
        build_slide_18(prs, visual_18)
        prs.save(args.output)
        configure_slide17_video(args.output)
        print(args.output)
        return

    if starting_slide_count == 18:
        build_slide_19(prs, visual_19)
        build_slide_20(prs, visual_20)
        build_slide_21(prs, visual_21)
        build_slide_23(prs, visual_23)
        prs.save(args.output)
        print(args.output)
        return

    if starting_slide_count == 22:
        build_slide_23_5(prs, visual_23_5)
        build_slide_24(prs, visual_24)
        build_slide_25(prs, anchor_images, args.mascot)
        prs.save(args.output)
        print(args.output)
        return

    if starting_slide_count == 25:
        build_slide_26(prs, visual_26)
        build_slide_27(prs, visual_27)
        build_slide_28(prs, video_28, poster_28)
        normalize_all_act_rails(prs)
        bring_speech_overlays_to_front(prs)
        configure_video_autoplay_loop_timing(prs, 28, "Slide28_Video")
        prs.save(args.output)
        externalize_linked_video(args.output, 28, video_28)
        print(args.output)
        return

    if starting_slide_count == 28:
        print(args.output)
        return

    # Act III is appended non-destructively to the current 11-slide deck. Do
    # not run any of the legacy normalization/replacement passes against Acts
    # 0-II; they already contain approved media, layout, and click timing.
    if starting_slide_count == 11:
        build_slide_12(prs, visual_12_01, visual_12_02)
        build_slide_13(prs, app_cam_paths, rigged_cam_paths)
        build_slide_14(prs, visual_14)
        prepare_speech_overlays(prs, args.mascot, args.serious_mascot)
        configure_speech_overlay_click_timing(prs, 14)
        prs.save(args.output)
        print(args.output)
        return

    if len(prs.slides) == 3:
        build_slide_4(prs, visual_4)
        build_slide_5(prs, visual_5)
        build_slide_6(prs, visual_6)
    else:
        replace_slide_5_placeholder(prs, visual_5)

    if len(prs.slides) == 6:
        build_slide_7(prs, visual_7)
        build_slide_8(prs, visual_8, poster_8)
        build_slide_9(prs, visual_9)
        build_slide_10(prs, visual_10, poster_10)
        build_slide_11(prs, visual_11)

    if len(prs.slides) == 11:
        build_slide_12(prs)
        build_slide_13(prs)
        build_slide_14(prs)

    replace_slide_visual(
        prs, 6, visual_6, 6.50, 1.48, 6.30, 4.98, "Slide6_Visual", dark=True
    )
    replace_slide_visual(
        prs, 7, visual_7, 6.48, 1.48, 6.32, 4.98, "Slide7_Visual", dark=True
    )
    replace_slide_visual(
        prs, 9, visual_9, 6.54, 1.48, 6.26, 4.98, "Slide9_Visual", dark=True
    )
    replace_slide_video(
        prs,
        10,
        visual_10,
        poster_10,
        1.75,
        1.48,
        6.18,
        4.82,
        "Slide10_Video",
        dark=False,
    )
    replace_slide_visual(
        prs, 11, visual_11, 6.48, 1.48, 6.32, 4.82, "Slide11_Visual", dark=True
    )
    normalize_slide_8_layout(prs, visual_8, poster_8)
    configure_video_click_timing(prs, 8, "Slide8_Video", 195566)
    configure_video_click_timing(
        prs,
        10,
        "Slide10_Video",
        72233,
        followup_zoom_shape_name="SpeechOverlay_Group",
    )
    normalize_slide_10_title(prs)
    normalize_all_act_rails(prs)
    prepare_speech_overlays(prs, args.mascot, args.serious_mascot)
    bring_speech_overlays_to_front(prs)
    prs.save(args.output)
    if not args.skip_animations:
        apply_speech_animations(args.output, args.serious_mascot)
    print(args.output)


if __name__ == "__main__":
    main()
